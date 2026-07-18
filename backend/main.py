import asyncio
import datetime
import json
import logging
import logging.handlers
import os
import re
import time

# Load .env from backend directory if present (never committed to git)
_env_path = os.path.join(os.path.dirname(__file__), ".env")
if os.path.isfile(_env_path):
    with open(_env_path) as _f:
        for _line in _f:
            _line = _line.strip()
            if _line and not _line.startswith("#") and "=" in _line:
                _k, _v = _line.split("=", 1)
                os.environ.setdefault(_k.strip(), _v.strip())
import shutil
import subprocess
import tempfile
import uuid
from collections import deque
from concurrent.futures import ThreadPoolExecutor
from contextlib import asynccontextmanager as _asynccontextmanager
from pathlib import Path
from typing import Optional, Dict

try:
    import psutil as _psutil_import  # noqa: F401 — availability check at top level
except ImportError:
    _psutil_import = None  # type: ignore

import edge_tts
import httpx
from fastapi import FastAPI, HTTPException, BackgroundTasks, UploadFile, Form, File, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, field_validator

try:
    from slowapi import Limiter, _rate_limit_exceeded_handler
    from slowapi.util import get_remote_address
    from slowapi.errors import RateLimitExceeded
    _SLOWAPI_AVAILABLE = True
except ImportError:
    _SLOWAPI_AVAILABLE = False

# ── Monkey-patch for coqui-tts / transformers compatibility ──────────────────
try:
    import torch as _torch
    import transformers.pytorch_utils as _tpu
    if not hasattr(_tpu, 'isin_mps_friendly'):
        _tpu.isin_mps_friendly = _torch.isin
except Exception:
    pass

try:
    # coqui-tts >= 0.27 raises ImportError when torch >= 2.9 and torchcodec is missing.
    # Patch is_torchcodec_available to bypass the startup check.
    import transformers.utils.import_utils as _triuu
    _triuu.is_torchcodec_available = lambda: True
except Exception:
    pass

try:
    # torchaudio 2.9+ delegates torchaudio.load() to torchcodec, which requires
    # GPU CUDA libraries we don't have. Replace with a soundfile-based implementation
    # that works for PCM WAV (which is all XTTS needs for reference audio).
    import torchaudio as _torchaudio
    import soundfile as _sf
    import torch as _torch_
    import numpy as _np

    def _sf_torchaudio_load(uri, frame_offset=0, num_frames=-1,
                            normalize=True, channels_first=True,
                            format=None, buffer_size=4096, backend=None):
        data, samplerate = _sf.read(str(uri), dtype='float32', always_2d=True,
                                    start=frame_offset,
                                    frames=num_frames if num_frames > 0 else -1)
        tensor = _torch_.from_numpy(_np.ascontiguousarray(data.T))  # (channels, samples)
        return tensor, samplerate

    _torchaudio.load = _sf_torchaudio_load
except Exception:
    pass

# ── Logging setup ────────────────────────────────────────────────────────────
logger = logging.getLogger("voice_studio")
logger.setLevel(os.getenv("LOG_LEVEL", "INFO").upper())
_log_fmt = logging.Formatter("%(asctime)s [%(levelname)s] %(name)s: %(message)s")
_sh = logging.StreamHandler()
_sh.setFormatter(_log_fmt)
logger.addHandler(_sh)
try:
    _rfh = logging.handlers.RotatingFileHandler(
        "/tmp/voice_studio.log", maxBytes=10 * 1024 * 1024, backupCount=3
    )
    _rfh.setFormatter(_log_fmt)
    logger.addHandler(_rfh)
except Exception:
    pass
logging.getLogger("uvicorn.access").setLevel(logging.WARNING)
logging.getLogger("httpx").setLevel(logging.WARNING)

# ── Job tracking (TTL dict — auto-expires entries after 2h) ──────────────────
class _TTLDict(dict):
    """dict subclass that tracks insertion timestamps and supports TTL cleanup.
    Synchronous interface (compatible with JOBS[k]=v usage throughout codebase).
    """
    def __init__(self, ttl_seconds: int = 7200):
        super().__init__()
        self._ttl = ttl_seconds
        self._timestamps: dict = {}

    def __setitem__(self, key, value):
        super().__setitem__(key, value)
        self._timestamps[key] = time.monotonic()

    def __delitem__(self, key):
        super().__delitem__(key)
        self._timestamps.pop(key, None)

    def cleanup(self) -> int:
        """Remove entries older than TTL. Returns number removed."""
        now = time.monotonic()
        expired = [k for k, ts in self._timestamps.items() if now - ts > self._ttl]
        for k in expired:
            super().pop(k, None)
            self._timestamps.pop(k, None)
        return len(expired)

JOBS: _TTLDict = _TTLDict(ttl_seconds=7200)

async def _jobs_cleanup_loop():
    """Periodic JOBS cleanup task — runs every 10 minutes."""
    while True:
        await asyncio.sleep(600)
        removed = JOBS.cleanup()
        if removed:
            logger.info("JOBS cleanup: removed %d expired entries", removed)

# XTTS model instance (loaded lazily, shared across requests)
_xtts_model = None
_xtts_lock = asyncio.Lock()
_xtts_executor = ThreadPoolExecutor(max_workers=1)

# Chatterbox TTS worker (separate venv — torch 2.6 vs our torch 2.8 can't coexist)
_chatterbox_proc: Optional[asyncio.subprocess.Process] = None
_chatterbox_start_lock = asyncio.Lock()
_chatterbox_rw_lock = asyncio.Lock()
_CHATTERBOX_VENV = Path(__file__).parent / ".chatterbox-venv"
_CHATTERBOX_WORKER = Path(__file__).parent / "chatterbox_worker.py"

# GPU availability (checked once at startup)
_GPU_AVAILABLE: Optional[bool] = None

def _check_gpu() -> bool:
    global _GPU_AVAILABLE
    if _GPU_AVAILABLE is not None:
        return _GPU_AVAILABLE
    try:
        import torch
        _GPU_AVAILABLE = torch.cuda.is_available()
    except Exception:
        _GPU_AVAILABLE = False
    return _GPU_AVAILABLE

def _load_xtts_sync(models_dir: str):
    """Synchronously load XTTS-v2 model (runs in thread executor). Uses GPU if available."""
    import os as _os
    _os.environ.setdefault("TTS_AGREE_TO_USER_AGREEMENT", "1")
    _os.environ.setdefault("COQUI_TOS_AGREED", "1")
    from TTS.api import TTS as CoquiTTS
    use_gpu = _check_gpu()
    try:
        tts = CoquiTTS("tts_models/multilingual/multi-dataset/xtts_v2", gpu=use_gpu)
        logger.info("[XTTS] Załadowano na %s", 'GPU' if use_gpu else 'CPU')
        return tts
    except Exception as e:
        if use_gpu:
            logger.warning("[XTTS] GPU load failed (%s), retry on CPU", e)
            return CoquiTTS("tts_models/multilingual/multi-dataset/xtts_v2", gpu=False)
        raise


# ── GPU Job Queue (FIFO, max 1 GPU job at a time) ──────────────────────────
class _GPUJobQueue:
    """Singleton async semaphore for GPU-bound jobs. Keeps VRAM contention safe."""
    def __init__(self):
        self._sem = asyncio.Semaphore(1)
        self._waiting: int = 0

    @_asynccontextmanager
    async def acquire(self):
        self._waiting += 1
        try:
            async with self._sem:
                self._waiting -= 1
                yield
        except Exception:
            self._waiting = max(0, self._waiting - 1)
            raise

    def status(self) -> dict:
        return {
            "active": self._sem._value == 0,
            "queue_length": self._waiting,
        }


GPU_QUEUE = _GPUJobQueue()

async def get_xtts_model():
    """Return cached XTTS model, loading on first call."""
    global _xtts_model
    if _xtts_model is not None:
        return _xtts_model
    async with _xtts_lock:
        if _xtts_model is None:
            loop = asyncio.get_event_loop()
            _xtts_model = await loop.run_in_executor(
                _xtts_executor, _load_xtts_sync, str(MODELS_DIR)
            )
    return _xtts_model

def _release_tts_vram() -> None:
    """Unload XTTS + stop Chatterbox worker before face animation.
    Both reload automatically on the next synthesis request (~30 s for XTTS,
    ~60 s for Chatterbox cold start). Frees 4–8 GB VRAM for EchoMimic/SadTalker."""
    global _xtts_model, _chatterbox_proc
    freed_something = False

    # 1. Release XTTS in-process
    if _xtts_model is not None:
        try:
            _xtts_model = None
            freed_something = True
        except Exception as _e:
            logger.debug("[VRAM] xtts release: %s", _e)
            _xtts_model = None

    # 2. Stop Chatterbox subprocess (holds 4+ GB VRAM in its own process)
    proc = _chatterbox_proc
    if proc is not None and proc.returncode is None:
        try:
            proc.terminate()
            freed_something = True
            logger.info("[VRAM] Chatterbox worker terminated — freeing ~4 GB VRAM for face animation")
        except Exception as _e:
            logger.debug("[VRAM] chatterbox terminate: %s", _e)
        _chatterbox_proc = None

    if freed_something:
        try:
            import torch
            import gc
            gc.collect()
            torch.cuda.empty_cache()
        except Exception:
            pass


_XTTS_CHAR_LIMIT = 200  # safe per-chunk limit for Polish XTTS (hard limit is 224)
_POLISH_CHARS = frozenset('ąęóśłźżćńĄĘÓŚŁŹŻĆŃ')
_XTTS_SAMPLE_RATE = 22050  # XTTS expected sample rate for reference audio
# XTTS-v2 generates language-conditioning artifacts (~200ms of CJK-sounding tokens) at the very
# start of every synthesis. Prepending a short warmup phrase causes those tokens to attach to the
# warmup instead of the real text, so trimming 450ms safely removes them without cutting real speech.
_XTTS_WARMUP = "Hm. "
_XTTS_TRIM_MS = 450


def _convert_to_xtts_wav(src: str, dst: str) -> None:
    """Convert any audio file to 22050Hz mono PCM WAV required by XTTS."""
    subprocess.run(
        ["ffmpeg", "-y", "-i", src, "-ar", str(_XTTS_SAMPLE_RATE), "-ac", "1",
         "-sample_fmt", "s16", "-f", "wav", dst],
        check=True, capture_output=True,
    )

def _detect_xtts_lang(text: str) -> str:
    """Detect language for XTTS: 'pl' if Polish diacritics found, else 'en'."""
    return 'pl' if any(c in _POLISH_CHARS for c in text) else 'en'


def _resolve_xtts_lang(text: str, voice_lang: Optional[str] = None) -> str:
    """Polish diacritics always force 'pl'; otherwise use stored voice language or 'en'."""
    if any(c in _POLISH_CHARS for c in text):
        return 'pl'
    return voice_lang or 'en'


def _save_voice_meta(safe_name: str, lang: str) -> None:
    import json as _json
    meta_path = XTTS_SPEAKERS_DIR / f"{safe_name}.json"
    meta_path.write_text(_json.dumps({"lang": lang}))


def _load_voice_lang(safe_name: str) -> Optional[str]:
    import json as _json
    meta_path = XTTS_SPEAKERS_DIR / f"{safe_name}.json"
    if meta_path.exists():
        try:
            return _json.loads(meta_path.read_text()).get("lang")
        except Exception:
            return None
    return None

def _sanitize_xtts_text(text: str) -> str:
    """Remove characters that cause XTTS to generate CJK artifacts or language switches.
    Keeps Polish/Latin letters, digits, standard punctuation."""
    import unicodedata
    allowed_punct = set('.,!?;:-—–()[]"\' \n\t')
    result = []
    for c in text:
        cp = ord(c)
        # Drop CJK, Arabic, Hebrew, Thai etc. — anything that could trigger lang switch
        if 0x2E80 <= cp <= 0x9FFF or 0xAC00 <= cp <= 0xD7FF or 0xF900 <= cp <= 0xFAFF:
            continue
        cat = unicodedata.category(c)
        if cat.startswith('L') or cat.startswith('N') or c in allowed_punct:
            result.append(c)
    text = ''.join(result)
    # Collapse whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    # Replace em-dashes mid-sentence with comma for more natural pacing
    text = text.replace('—', ',').replace('–', ',')
    # Exclamation mark at end reduces XTTS hallucinations more reliably than period (community finding)
    if text and text[-1] not in '!?…':
        text = text.rstrip('.') + '!'
    return text


def _trim_xtts_wav(wav_path: str, trim_ms: int = 250) -> None:
    """Trim first N ms from XTTS output to eliminate language-artifact tokens (Chinese etc.)."""
    tmp = wav_path + ".trimmed.wav"
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", wav_path, "-ss", f"{trim_ms / 1000:.3f}", tmp],
            check=True, capture_output=True,
        )
        os.replace(tmp, wav_path)
    except Exception:
        try: os.unlink(tmp)
        except OSError: pass


def _trim_xtts_wav_tail(wav_path: str) -> None:
    """Remove trailing silence + hallucinated tokens from XTTS output.
    Uses areverse trick: reverse audio → remove leading silence → reverse back → add 120ms clean pad."""
    tmp = wav_path + ".tail.wav"
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", wav_path,
             "-af", "areverse,silenceremove=start_periods=1:start_duration=0.12:start_threshold=-38dB:detection=rms,areverse,apad=pad_dur=0.12",
             tmp],
            check=True, capture_output=True,
        )
        os.replace(tmp, wav_path)
    except Exception:
        try: os.unlink(tmp)
        except OSError: pass


def _xtts_synthesize_sync(tts, text: str, speaker_wav: str, out_path: str, language: str = "pl", speed: float = 1.0):
    text = _sanitize_xtts_text(text)
    if not text:
        return
    # Prepend warmup so XTTS conditioning tokens attach to it, not to the real text.
    tts.tts_to_file(
        text=_XTTS_WARMUP + text,
        speaker_wav=speaker_wav,
        language=language,
        file_path=out_path,
        speed=speed,
        temperature=0.55,
        repetition_penalty=3.0,
        top_k=40,
        top_p=0.75,
    )
    # Trim warmup phrase + conditioning tokens (covers ~150ms CJK tokens + ~200ms "Hm." speech)
    _trim_xtts_wav(out_path, trim_ms=_XTTS_TRIM_MS)

def _split_for_xtts(text: str) -> list:
    """Split text into XTTS-safe chunks (max _XTTS_CHAR_LIMIT chars) at sentence/clause boundaries."""
    import re
    if len(text) <= _XTTS_CHAR_LIMIT:
        return [text]
    parts = re.split(r'(?<=[.!?…])\s+', text)
    chunks, current = [], ""
    for part in parts:
        if len(current) + len(part) + 1 <= _XTTS_CHAR_LIMIT:
            current = (current + " " + part).strip() if current else part
        else:
            if current:
                chunks.append(current)
            if len(part) > _XTTS_CHAR_LIMIT:
                for sub in re.split(r'(?<=[,;])\s+', part):
                    if len(sub) > _XTTS_CHAR_LIMIT:
                        # Last resort: split at word boundaries, never mid-word
                        word_chunk = ""
                        for word in sub.split():
                            if len(word_chunk) + len(word) + 1 <= _XTTS_CHAR_LIMIT:
                                word_chunk = (word_chunk + " " + word).strip() if word_chunk else word
                            else:
                                if word_chunk:
                                    chunks.append(word_chunk)
                                word_chunk = word
                        if word_chunk:
                            chunks.append(word_chunk)
                    else:
                        chunks.append(sub)
                current = ""
            else:
                current = part
    if current:
        chunks.append(current)
    return [c for c in chunks if c.strip()]

def _ffmpeg_concat_wavs(input_paths: list, out_path: str) -> None:
    """Concatenate WAV files using ffmpeg (no pydub needed)."""
    list_file = out_path + ".concat.txt"
    try:
        with open(list_file, "w") as f:
            for p in input_paths:
                f.write(f"file '{p}'\n")
        subprocess.run(
            ["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", list_file, out_path],
            check=True, capture_output=True,
        )
    finally:
        try:
            os.unlink(list_file)
        except OSError:
            pass


def _ffmpeg_wav_to_mp3(wav_path: str, mp3_path: str) -> None:
    """Convert WAV to MP3 using ffmpeg (no pydub needed)."""
    subprocess.run(
        ["ffmpeg", "-y", "-i", wav_path, "-q:a", "4", mp3_path],
        check=True, capture_output=True,
    )


def _ffmpeg_mp3_duration_ms(mp3_path: str) -> int:
    """Return duration of an MP3 file in milliseconds using ffprobe."""
    import json as _json
    r = subprocess.run(
        ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", mp3_path],
        capture_output=True, text=True, check=True,
    )
    return int(float(_json.loads(r.stdout)["format"]["duration"]) * 1000)


async def xtts_synthesize(text: str, speaker_wav: str, out_path: str, speed: float = 1.0, force_language: Optional[str] = None):
    """Async wrapper: synthesize voice clone via local XTTS. Splits long text, resolves language."""
    import tempfile
    chunks = _split_for_xtts(text)
    tts = await get_xtts_model()
    loop = asyncio.get_event_loop()
    if len(chunks) == 1:
        language = _resolve_xtts_lang(chunks[0], force_language)
        await loop.run_in_executor(
            _xtts_executor, _xtts_synthesize_sync, tts, chunks[0], speaker_wav, out_path, language, speed
        )
        return
    tmp_paths = []
    try:
        for chunk in chunks:
            chunk_lang = _resolve_xtts_lang(chunk, force_language)
            tf = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
            tf.close()
            await loop.run_in_executor(
                _xtts_executor, _xtts_synthesize_sync, tts, chunk, speaker_wav, tf.name, chunk_lang, speed
            )
            tmp_paths.append(tf.name)
        _ffmpeg_concat_wavs(tmp_paths, out_path)
    finally:
        for p in tmp_paths:
            try:
                os.unlink(p)
            except OSError:
                pass

async def _get_chatterbox_proc() -> Optional[asyncio.subprocess.Process]:
    """Start (or reuse) the Chatterbox worker process. Returns None if not installed."""
    global _chatterbox_proc
    py = _CHATTERBOX_VENV / "bin" / "python"
    if not py.exists() or not _CHATTERBOX_WORKER.exists():
        return None
    if _chatterbox_proc is not None and _chatterbox_proc.returncode is None:
        return _chatterbox_proc
    async with _chatterbox_start_lock:
        if _chatterbox_proc is None or _chatterbox_proc.returncode is not None:
            logger.info("Starting Chatterbox worker (first use — model loading ~30s)…")
            _chatterbox_proc = await asyncio.create_subprocess_exec(
                str(py), str(_CHATTERBOX_WORKER),
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.DEVNULL,
            )
            ready = await asyncio.wait_for(_chatterbox_proc.stdout.readline(), timeout=180)
            if not json.loads(ready).get("ready"):
                raise RuntimeError("Chatterbox worker failed to start")
            logger.info("Chatterbox worker ready")
    return _chatterbox_proc


async def _chatterbox_synthesize_chunk(proc, text: str, ref_wav: Optional[str], out_wav: str, language: str) -> None:
    """Send one chunk to the Chatterbox worker and wait for the WAV response."""
    req = json.dumps({"text": text, "ref_wav": ref_wav, "out_wav": out_wav, "lang": language})
    async with _chatterbox_rw_lock:
        proc.stdin.write((req + "\n").encode())
        await proc.stdin.drain()
        resp = await asyncio.wait_for(proc.stdout.readline(), timeout=120)
    result = json.loads(resp)
    if not result.get("ok"):
        raise RuntimeError(f"Chatterbox: {result.get('error', 'unknown')}")


async def chatterbox_synthesize(text: str, ref_wav: Optional[str], out_path_mp3: str, language: str = "pl") -> None:
    """Synthesize via Chatterbox worker (EN=Turbo, other=Multilingual V3). Output: MP3.

    Splits long text into ~200-char sentence chunks to avoid hitting the model's
    generation step limit (~1000 steps ≈ 30s), which causes mid-sentence cutoffs.
    """
    _CHUNK_LIMIT = 200
    proc = await _get_chatterbox_proc()
    if proc is None:
        raise RuntimeError("Chatterbox not installed — run install.sh")

    chunks = _split_for_xtts(text) if len(text) > _CHUNK_LIMIT else [text]
    tmp_wavs: list[str] = []
    try:
        for idx, chunk in enumerate(chunks):
            wav_path = out_path_mp3 + f".chtrbx_{idx}.wav"
            tmp_wavs.append(wav_path)
            await _chatterbox_synthesize_chunk(proc, chunk, ref_wav, wav_path, language)

        if len(tmp_wavs) == 1:
            merged_wav = tmp_wavs[0]
        else:
            merged_wav = out_path_mp3 + ".chtrbx_merged.wav"
            tmp_wavs.append(merged_wav)
            _ffmpeg_concat_wavs(tmp_wavs[:-1], merged_wav)

        conv = await asyncio.create_subprocess_exec(
            "ffmpeg", "-y", "-i", merged_wav, "-q:a", "4", out_path_mp3,
            stderr=asyncio.subprocess.DEVNULL, stdout=asyncio.subprocess.DEVNULL,
        )
        await conv.communicate()
    finally:
        for p in tmp_wavs:
            try:
                os.unlink(p)
            except OSError:
                pass


STYLE_MAP: Dict[str, str] = {
    'cinematic': 'cinematic atmosphere',
    'anime': 'anime style',
    'realistic': 'realistic',
    '3d_render': '3D rendered',
    'cyberpunk': 'cyberpunk neon',
    'horror': 'dark horror atmosphere',
}

# ── Simple rate limiter fallback (used when slowapi not installed) ─────────────
class _SimpleRateLimiter:
    def __init__(self):
        self._history: dict = {}
        self._lock = asyncio.Lock()

    def _parse_rule(self, rule: str):
        count, period = rule.split("/")
        return int(count), {"second": 1, "minute": 60, "hour": 3600}.get(period.lower(), 60)

    async def check(self, key: str, rule: str):
        limit, window = self._parse_rule(rule)
        now = time.monotonic()
        async with self._lock:
            history = [t for t in self._history.get(key, []) if now - t < window]
            if len(history) >= limit:
                raise HTTPException(429, detail="Too Many Requests. Spróbuj za chwilę.")
            history.append(now)
            self._history[key] = history

    async def cleanup_loop(self):
        while True:
            await asyncio.sleep(60)
            now = time.monotonic()
            async with self._lock:
                for key in list(self._history.keys()):
                    self._history[key] = [t for t in self._history[key] if now - t < 3600]
                    if not self._history[key]:
                        del self._history[key]

_simple_limiter = _SimpleRateLimiter()

# ── Lifespan (startup/shutdown) ───────────────────────────────────────────────
@_asynccontextmanager
async def lifespan(app: FastAPI):
    asyncio.create_task(_jobs_cleanup_loop())
    if not _SLOWAPI_AVAILABLE:
        asyncio.create_task(_simple_limiter.cleanup_loop())
    # Video AI Orchestrator
    try:
        from video_ai import init_db as _va_init_db, start_worker as _va_start_worker
        _va_init_db()
        asyncio.create_task(_va_start_worker())
        logger.info("Video AI Orchestrator started")
    except Exception as _e:
        logger.warning("Video AI Orchestrator not started: %s", _e)
    logger.info("Voice Studio startup complete")
    yield
    logger.info("Voice Studio shutting down")

app = FastAPI(title="Voice & Video Studio AI", version="2.0.0", lifespan=lifespan)

# ── Rate Limiting ─────────────────────────────────────────────────────────────
if _SLOWAPI_AVAILABLE:
    _limiter = Limiter(key_func=get_remote_address, default_limits=["200/minute"])
    app.state.limiter = _limiter
    app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)
else:
    _limiter = None

def _rate_limit(rule: str = "100/minute"):
    """Decorator that applies rate limit only if slowapi is installed."""
    def decorator(func):
        if _SLOWAPI_AVAILABLE and _limiter:
            return _limiter.limit(rule)(func)
        return func
    return decorator

# ── Job Queue System ──────────────────────────────────────────────────────────
_JOB_QUEUE: deque = deque()
_JOB_QUEUE_RUNNING: bool = False
_JOB_QUEUE_ACTIVE: Optional[str] = None
_JOB_QUEUE_LOCK = asyncio.Lock()

async def _queue_worker():
    global _JOB_QUEUE_RUNNING, _JOB_QUEUE_ACTIVE
    _JOB_QUEUE_RUNNING = True
    while _JOB_QUEUE:
        job_id, func, args, kwargs = _JOB_QUEUE.popleft()
        _JOB_QUEUE_ACTIVE = job_id
        if job_id in JOBS:
            JOBS[job_id].update({"status": "running", "queue_position": None, "message": "Uruchamianie…"})
        try:
            await func(*args, **kwargs)
        except Exception as e:
            if job_id in JOBS:
                JOBS[job_id] = {"status": "failed", "error": str(e), "url": None, "progress": None}
        _JOB_QUEUE_ACTIVE = None
    _JOB_QUEUE_RUNNING = False

async def _enqueue_job(job_id: str, func, *args, **kwargs):
    global _JOB_QUEUE_RUNNING
    async with _JOB_QUEUE_LOCK:
        _JOB_QUEUE.append((job_id, func, args, kwargs))
        JOBS[job_id] = {
            "status": "queued", "progress": 0, "url": None, "error": None,
            "queue_position": len(_JOB_QUEUE),
            "created_ts": datetime.datetime.now().isoformat(),
        }
        for i, (jid, *_) in enumerate(_JOB_QUEUE):
            if jid in JOBS:
                JOBS[jid]["queue_position"] = i + 1
        if not _JOB_QUEUE_RUNNING:
            _JOB_QUEUE_RUNNING = True
            asyncio.create_task(_queue_worker())

# ── Night Mode ────────────────────────────────────────────────────────────────
_NIGHT_MODE: dict = {
    "enabled": False,
    "start_hour": 22,
    "end_hour": 7,
}

def _is_night_mode_active() -> bool:
    if not _NIGHT_MODE["enabled"]:
        return False
    hour = datetime.datetime.now().hour
    s, e = _NIGHT_MODE["start_hour"], _NIGHT_MODE["end_hour"]
    if s > e:  # spans midnight e.g. 22→07
        return hour >= s or hour < e
    return s <= hour < e

def _night_mode_block():
    if _is_night_mode_active():
        s, e = _NIGHT_MODE["start_hour"], _NIGHT_MODE["end_hour"]
        raise HTTPException(503, detail={
            "error": "night_mode_active",
            "message": f"Tryb nocny aktywny ({s:02d}:00–{e:02d}:00). Zadania GPU wstrzymane do rana.",
            "retry_after": f"{e:02d}:00",
        })

_cors_origins_env = os.getenv("CORS_ORIGINS", "")
_cors_origins = (
    [o.strip() for o in _cors_origins_env.split(",") if o.strip()]
    if _cors_origins_env
    else ["http://localhost:47822", "http://127.0.0.1:47822"]
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins,
    allow_credentials=True,
    allow_methods=["GET", "POST", "DELETE"],
    allow_headers=["Content-Type", "Authorization"],
)

# ── Video AI Orchestrator router ──────────────────────────────────────────────
try:
    from video_ai import router as _video_ai_router
    app.include_router(_video_ai_router)
    logger.info("Video AI router registered")
except Exception as _e:
    logger.warning("Video AI router not loaded: %s", _e)

# PATH CONFIGURATION
# Set VOICE_STUDIO_DATA_DIR env var to use an external drive / custom path.
# Falls back to <project_root>/audio-output, ai-models, xtts-speakers.
_EXTERNAL_ROOT = Path(os.environ.get("VOICE_STUDIO_DATA_DIR", "")) if os.environ.get("VOICE_STUDIO_DATA_DIR") else None
_LOCAL_ROOT = Path(__file__).parent.parent

_AUDIO_EXT  = _EXTERNAL_ROOT / "voice-studio-audio" if _EXTERNAL_ROOT else None
_LOCAL_AUDIO = _LOCAL_ROOT / "audio-output"
AUDIO_DIR = _AUDIO_EXT if (_AUDIO_EXT and _EXTERNAL_ROOT.exists()) else _LOCAL_AUDIO
AUDIO_DIR.mkdir(parents=True, exist_ok=True)

_MODELS_EXT  = _EXTERNAL_ROOT / "ai-models" if _EXTERNAL_ROOT else None
_LOCAL_MODELS = _LOCAL_ROOT / "ai-models"
MODELS_DIR = _MODELS_EXT if (_MODELS_EXT and _EXTERNAL_ROOT.exists()) else _LOCAL_MODELS
MODELS_DIR.mkdir(parents=True, exist_ok=True)

_SPEAKERS_EXT  = _EXTERNAL_ROOT / "xtts-speakers" if _EXTERNAL_ROOT else None
_LOCAL_SPEAKERS = _LOCAL_ROOT / "xtts-speakers"
XTTS_SPEAKERS_DIR = _SPEAKERS_EXT if (_SPEAKERS_EXT and _EXTERNAL_ROOT.exists()) else _LOCAL_SPEAKERS
XTTS_SPEAKERS_DIR.mkdir(parents=True, exist_ok=True)

# Serve built frontend if dist exists
FRONTEND_DIST = _LOCAL_ROOT / "frontend" / "dist"

# ──────────────────────────────────────────────────────────────────────────────
# VOICE CATALOG
# ──────────────────────────────────────────────────────────────────────────────
VOICES = {
    # POLISH (Native)
    "pl_male_marek": {
        "id": "pl-PL-MarekNeural", "lang": "pl-PL",
        "gender": "male", "age": "adult", "label": "Marek (PL, Mężczyzna)",
        "personality": "Poważny, profesjonalny"
    },
    "pl_female_zofia": {
        "id": "pl-PL-ZofiaNeural", "lang": "pl-PL",
        "gender": "female", "age": "adult", "label": "Zofia (PL, Kobieta)",
        "personality": "Przyjazna, ciepła"
    },
    # MULTILINGUAL (Excellent Polish support)
    "pl_male_andrew_multi": {
        "id": "en-US-AndrewMultilingualNeural", "lang": "pl-PL",
        "gender": "male", "age": "adult", "label": "Andrew (Multilingual, M)",
        "personality": "Płynny, ciepły, bardzo naturalny"
    },
    "pl_female_ava_multi": {
        "id": "en-US-AvaMultilingualNeural", "lang": "pl-PL",
        "gender": "female", "age": "adult", "label": "Ava (Multilingual, K)",
        "personality": "Ekspresyjna, miła"
    },
    "pl_male_brian_multi": {
        "id": "en-US-BrianMultilingualNeural", "lang": "pl-PL",
        "gender": "male", "age": "adult", "label": "Brian (Multilingual, M)",
        "personality": "Naturalny, codzienny"
    },
    "pl_female_emma_multi": {
        "id": "en-US-EmmaMultilingualNeural", "lang": "pl-PL",
        "gender": "female", "age": "adult", "label": "Emma (Multilingual, K)",
        "personality": "Radosna, czysta"
    },
    "pl_female_vivienne_multi": {
        "id": "fr-FR-VivienneMultilingualNeural", "lang": "pl-PL",
        "gender": "female", "age": "adult", "label": "Vivienne (Multilingual, K)",
        "personality": "Elegancka, spokojna"
    },
    "pl_female_seraphina_multi": {
        "id": "de-DE-SeraphinaMultilingualNeural", "lang": "pl-PL",
        "gender": "female", "age": "adult", "label": "Seraphina (Multilingual, K)",
        "personality": "Precyzyjna, inteligentna"
    },
    "pl_male_giuseppe_multi": {
        "id": "it-IT-GiuseppeMultilingualNeural", "lang": "pl-PL",
        "gender": "male", "age": "adult", "label": "Giuseppe (Multilingual, M)",
        "personality": "Energiczny, pasjonujący"
    },
    # ENGLISH MALE
    "en_male_guy": {
        "id": "en-US-GuyNeural", "lang": "en-US",
        "gender": "male", "age": "adult", "label": "Guy (EN, man)",
        "personality": "Energetic, Passionate"
    },
    "en_male_andrew": {
        "id": "en-US-AndrewNeural", "lang": "en-US",
        "gender": "male", "age": "adult", "label": "Andrew (EN, man)",
        "personality": "Warm, Confident"
    },
    "en_male_eric": {
        "id": "en-US-EricNeural", "lang": "en-US",
        "gender": "male", "age": "adult", "label": "Eric (EN, man)",
        "personality": "Rational, Calm"
    },
    "en_male_christopher": {
        "id": "en-US-ChristopherNeural", "lang": "en-US",
        "gender": "male", "age": "adult", "label": "Christopher (EN, man)",
        "personality": "Authoritative, Reliable"
    },
    # ENGLISH FEMALE
    "en_female_aria": {
        "id": "en-US-AriaNeural", "lang": "en-US",
        "gender": "female", "age": "adult", "label": "Aria (EN, woman)",
        "personality": "Positive, Confident"
    },
    "en_female_jenny": {
        "id": "en-US-JennyNeural", "lang": "en-US",
        "gender": "female", "age": "adult", "label": "Jenny (EN, woman)",
        "personality": "Friendly, Considerate"
    },
    # CHILD
    "en_child_ana": {
        "id": "en-US-AnaNeural", "lang": "en-US",
        "gender": "female", "age": "child", "label": "Ana (EN, girl)",
        "personality": "Cute, Playful"
    },
    "en_child_maisie_gb": {
        "id": "en-GB-MaisieNeural", "lang": "en-GB",
        "gender": "female", "age": "child", "label": "Maisie (EN-GB, girl)",
        "personality": "British, Cheerful"
    }
}

# ──────────────────────────────────────────────────────────────────────────────
# MODELS
# ──────────────────────────────────────────────────────────────────────────────

class SpeakerSegment(BaseModel):
    speaker_key: str        # key from VOICES dict
    text: str
    rate: Optional[str] = "+0%"   # e.g. "-10%", "+20%"  (edge-tts)
    pitch: Optional[str] = "+0Hz"  # e.g. "-20Hz", "+10Hz" (edge-tts)
    volume: Optional[str] = "+0%"
    xtts_speed: Optional[float] = 1.0  # XTTS speaking speed (0.5–2.0, 1.0 = match voice sample)

class GenerateRequest(BaseModel):
    segments: list[SpeakerSegment]
    silence_between_ms: Optional[int] = 500
    style: Optional[str] = "normal"   # normal | dramatic | calm | cheerful | newscast
    output_format: Optional[str] = "mp3"

class AIRequest(BaseModel):
    prompt: str
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"

# Legacy aliases
PromptRequest = AIRequest
ScriptRequest = AIRequest

class OllamaConfig(BaseModel):
    host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"

# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

STYLE_PROSODY = {
    "normal":    {"rate": "+0%",  "pitch": "+0Hz"},
    "dramatic":  {"rate": "-10%", "pitch": "-15Hz"},
    "calm":      {"rate": "-15%", "pitch": "-5Hz"},
    "cheerful":  {"rate": "+15%", "pitch": "+10Hz"},
    "newscast":  {"rate": "+5%",  "pitch": "+0Hz"},
    "whisper":   {"rate": "-20%", "pitch": "-30Hz"},
    "excited":   {"rate": "+25%", "pitch": "+20Hz"},
}

async def generate_segment_audio(segment: SpeakerSegment, style: str, tmp_dir: str) -> str:
    """Generate audio for a single segment, return file path."""
    # We need to re-fetch catalog to include dynamic cloned voices
    full_catalog = await get_full_catalog()
    voice_info = full_catalog.get(segment.speaker_key)
    
    if not voice_info:
        raise HTTPException(status_code=400, detail=f"Unknown speaker: {segment.speaker_key}")

    out_path = os.path.join(tmp_dir, f"{uuid.uuid4().hex}.mp3")

    # CASE A: Cloned Voice (XTTS-v2 via local model, edge-tts fallback)
    if voice_info.get("type") == "cloned":
        speaker_wav = str(XTTS_SPEAKERS_DIR / voice_info["id"])
        try:
            wav_path = out_path.replace(".mp3", ".wav")
            voice_name = segment.speaker_key[len('cloned_'):]
            voice_lang = voice_info.get("voice_lang") or _load_voice_lang(voice_name)
            await xtts_synthesize(segment.text, speaker_wav, wav_path, speed=segment.xtts_speed or 1.0, force_language=voice_lang)
            _ffmpeg_wav_to_mp3(wav_path, out_path)
            if os.path.exists(wav_path):
                os.unlink(wav_path)
            return out_path
        except Exception as e:
            import traceback as _tb
            logger.exception("[XTTS] Błąd syntezy (%s), fallback → edge-tts", e)
            fallback_voice = "pl-PL-MarekNeural"
            communicate = edge_tts.Communicate(
                text=segment.text, voice=fallback_voice, rate="+0%", pitch="+0Hz", volume="+0%"
            )
            await communicate.save(out_path)
            return out_path

    # CASE B: Standard Voice (edge-tts)
    voice_id = voice_info["id"]
    style_defaults = STYLE_PROSODY.get(style, STYLE_PROSODY["normal"])
    rate = segment.rate if segment.rate != "+0%" else style_defaults["rate"]
    pitch = segment.pitch if segment.pitch != "+0Hz" else style_defaults["pitch"]
    volume = segment.volume or "+0%"

    communicate = edge_tts.Communicate(
        text=segment.text,
        voice=voice_id,
        rate=rate,
        pitch=pitch,
        volume=volume,
    )
    await communicate.save(out_path)
    return out_path


FFMPEG_AVAILABLE: Optional[bool] = None

async def _check_ffmpeg() -> bool:
    global FFMPEG_AVAILABLE
    if FFMPEG_AVAILABLE is not None:
        return FFMPEG_AVAILABLE
    try:
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-version",
            stdout=asyncio.subprocess.DEVNULL,
            stderr=asyncio.subprocess.DEVNULL,
        )
        await proc.communicate()
        FFMPEG_AVAILABLE = proc.returncode == 0
    except Exception:
        FFMPEG_AVAILABLE = False
    return FFMPEG_AVAILABLE


async def create_silence(ms: int, tmp_dir: str) -> str:
    silence_path = os.path.join(tmp_dir, f"silence_{uuid.uuid4().hex}.mp3")
    if await _check_ffmpeg():
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-f", "lavfi", "-i", "anullsrc=r=24000:cl=mono",
            "-t", str(ms / 1000), "-q:a", "9", "-acodec", "libmp3lame",
            silence_path, "-y",
            stdout=asyncio.subprocess.DEVNULL,
            stderr=asyncio.subprocess.DEVNULL,
        )
        await proc.communicate()
        if proc.returncode == 0 and os.path.exists(silence_path):
            return silence_path
    subprocess.run(
        ["ffmpeg", "-f", "lavfi", "-i", "anullsrc=r=24000:cl=mono",
         "-t", str(ms / 1000), "-q:a", "9", "-acodec", "libmp3lame", silence_path, "-y"],
        check=True, capture_output=True,
    )
    return silence_path


async def _audiobook_master_chapter(src: str, dst: str) -> None:
    """Audiobook mastering chain: declick → compress → EQ → loudnorm → 44.1kHz/128kbps."""
    af_chain = ",".join([
        "highpass=f=80",
        "adeclick=w=55:o=75",
        "acompressor=threshold=-20dB:ratio=2.5:attack=15:release=250",
        "equalizer=f=120:width_type=q:width=1:g=2",
        "equalizer=f=8000:width_type=q:width=2:g=-1.5",
        "loudnorm=I=-18:TP=-3.0:LRA=7:print_format=none",
        "aresample=44100",
    ])
    proc = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", src, "-af", af_chain,
        "-c:a", "libmp3lame", "-b:a", "128k", "-ac", "1", dst,
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await proc.communicate()
    if not os.path.exists(dst):
        shutil.copy(src, dst)


async def concatenate_audio(files: list[str], out_path: str):
    if len(files) == 1:
        shutil.copy(files[0], out_path)
        return

    if await _check_ffmpeg():
        list_file = out_path + ".txt"
        try:
            with open(list_file, "w") as f:
                for fp in files:
                    f.write(f"file '{fp}'\n")
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-f", "concat", "-safe", "0", "-i", list_file,
                "-c", "copy", out_path, "-y",
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL,
            )
            await proc.communicate()
        finally:
            if os.path.exists(list_file):
                os.unlink(list_file)
        if os.path.exists(out_path):
            return

    try:
        _ffmpeg_concat_wavs(files, out_path)
    except Exception as e:
        if os.path.exists(out_path):
            os.unlink(out_path)
        raise RuntimeError(f"Błąd łączenia audio: {e}")

async def get_full_catalog():
    """Returns VOICES catalog merged with dynamically discovered cloned voices."""
    catalog = VOICES.copy()
    if XTTS_SPEAKERS_DIR.exists():
        for f in XTTS_SPEAKERS_DIR.glob("*"):
            if f.suffix.lower() in (".wav", ".mp3"):
                voice_key = f"cloned_{f.stem}"
                stored_lang = _load_voice_lang(f.stem) or 'pl'
                lang_label = 'EN' if stored_lang == 'en' else 'PL'
                catalog[voice_key] = {
                    "id": f.name,
                    "lang": stored_lang,
                    "gender": "unknown",
                    "age": "adult",
                    "label": f"👤 {f.stem} [{lang_label}] (Klon)",
                    "personality": "Głos sklonowany",
                    "type": "cloned",
                    "voice_lang": stored_lang,
                }
    return catalog

# ──────────────────────────────────────────────────────────────────────────────
# ROUTES
# ──────────────────────────────────────────────────────────────────────────────

@app.get("/api/voices")
async def get_voices():
    return {"voices": await get_full_catalog()}

@app.get("/api/styles")
async def get_styles():
    return {"styles": list(STYLE_PROSODY.keys())}

async def _process_generate_job(job_id: str, req: "GenerateRequest"):
    tmp_dir = tempfile.mkdtemp(prefix="vs_")
    segment_files = []
    total = len(req.segments)
    try:
        for i, seg in enumerate(req.segments):
            JOBS[job_id].update({
                "progress": int(10 + 80 * i / total),
                "message": f"Segment {i+1}/{total}…",
            })
            audio_path = await generate_segment_audio(seg, req.style or "normal", tmp_dir)
            segment_files.append(audio_path)
            if i < total - 1 and req.silence_between_ms and req.silence_between_ms > 0:
                silence_path = await create_silence(req.silence_between_ms, tmp_dir)
                segment_files.append(silence_path)

        JOBS[job_id].update({"progress": 95, "message": "Łączenie segmentów…"})
        output_filename = f"{job_id}.mp3"
        out_path = str(AUDIO_DIR / output_filename)
        await concatenate_audio(segment_files, out_path)
        JOBS[job_id].update({
            "status": "completed", "progress": 100, "message": "Gotowe!",
            "url": f"/api/audio/{output_filename}",
            "filename": f"voice-studio-{job_id[:8]}.mp3",
        })
    except Exception as e:
        JOBS[job_id].update({"status": "failed", "error": str(e)})
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@app.post("/api/generate")
async def generate_audio(req: GenerateRequest, background_tasks: BackgroundTasks):
    if not req.segments:
        raise HTTPException(status_code=400, detail="No segments provided")

    has_cloned = any(seg.speaker_key.startswith("cloned_") for seg in req.segments)

    if has_cloned:
        # Cloned voice uses Chatterbox/XTTS — must go through GPU queue to avoid
        # VRAM conflicts with concurrent presentation or audiobook rendering.
        job_id = uuid.uuid4().hex[:8]
        JOBS[job_id] = {"status": "queued", "progress": 0, "message": "Oczekuje w kolejce…", "mode": "voice_clone"}
        background_tasks.add_task(_enqueue_job, job_id, _process_generate_job, job_id, req)
        return {"job_id": job_id, "queued": True, "queue_length": len(_JOB_QUEUE) + 1}

    # Fast path: edge-tts only — synchronous
    job_id = uuid.uuid4().hex
    tmp_dir = tempfile.mkdtemp(prefix="vs_")
    segment_files = []
    try:
        for i, seg in enumerate(req.segments):
            audio_path = await generate_segment_audio(seg, req.style or "normal", tmp_dir)
            segment_files.append(audio_path)
            if i < len(req.segments) - 1 and req.silence_between_ms and req.silence_between_ms > 0:
                silence_path = await create_silence(req.silence_between_ms, tmp_dir)
                segment_files.append(silence_path)
        out_path = str(AUDIO_DIR / f"{job_id}.mp3")
        await concatenate_audio(segment_files, out_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    return {
        "job_id": job_id,
        "url": f"/api/audio/{job_id}.mp3",
        "filename": f"voice-studio-{job_id[:8]}.mp3"
    }

_AUDIO_SUFFIXES = {".mp3", ".mp4", ".wav", ".webm", ".m4b"}
_VIDEO_SUFFIXES = {".mp4", ".webm"}


def _sanitize_filename(name: str) -> str:
    import re
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '', name)
    name = re.sub(r'\s+', ' ', name).strip()
    return name[:80] or "audiobook"


@app.get("/api/files")
async def list_output_files():
    files, folders = [], []
    audio_dir_resolved = AUDIO_DIR.resolve()
    _INTERNAL_DIR_PREFIXES = ("movie-", "pres-", "ab-", "vgen-", "dialogue-")
    for entry in sorted(AUDIO_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
        if entry.is_dir() and any(entry.name.startswith(p) for p in _INTERNAL_DIR_PREFIXES):
            continue
        if entry.is_dir():
            folder_files = []
            for ff in sorted(entry.iterdir(), key=lambda x: x.stat().st_mtime):
                if ff.is_file() and ff.suffix.lower() in _AUDIO_SUFFIXES:
                    stat = ff.stat()
                    rel = f"{entry.name}/{ff.name}"
                    folder_files.append({
                        "name": ff.name,
                        "path": rel,
                        "size": stat.st_size,
                        "mtime": stat.st_mtime,
                        "type": "video" if ff.suffix.lower() in _VIDEO_SUFFIXES else "audio",
                        "url": f"/api/audio/{rel}",
                    })
            if folder_files:
                folders.append({
                    "name": entry.name,
                    "mtime": entry.stat().st_mtime,
                    "files": folder_files,
                })
        elif entry.is_file() and entry.suffix.lower() in _AUDIO_SUFFIXES:
            stat = entry.stat()
            files.append({
                "name": entry.name,
                "path": entry.name,
                "size": stat.st_size,
                "mtime": stat.st_mtime,
                "type": "video" if entry.suffix.lower() in _VIDEO_SUFFIXES else "audio",
                "url": f"/api/audio/{entry.name}",
            })
    return {"files": files, "folders": folders}




@app.post("/api/files/download-zip")
async def download_files_zip(request: Request):
    import zipfile, io as _io
    body = await request.json()
    paths = body.get("paths", [])
    buf = _io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in paths:
            try:
                target = _safe_audio_path(p)
                if target.is_file():
                    zf.write(target, target.name)
                elif target.is_dir():
                    for ff in sorted(target.iterdir()):
                        if ff.is_file() and ff.suffix.lower() in _AUDIO_SUFFIXES:
                            zf.write(ff, f"{target.name}/{ff.name}")
            except Exception:
                pass
    buf.seek(0)
    from fastapi.responses import StreamingResponse
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": "attachment; filename=projekty.zip"}
    )

def _safe_audio_path(filepath: str) -> Path:
    """Resolve filepath under AUDIO_DIR. Raises 403 on path traversal (including symlinks)."""
    try:
        base = AUDIO_DIR.resolve()
        target = (base / filepath).resolve()
        target.relative_to(base)  # ValueError if outside base
        return target
    except (ValueError, RuntimeError):
        raise HTTPException(status_code=403, detail="Access denied: niedozwolona ścieżka.")


@app.delete("/api/files")
async def delete_all_output_files():
    """Delete every file and folder in AUDIO_DIR."""
    deleted = 0
    for item in AUDIO_DIR.iterdir():
        try:
            if item.is_dir():
                shutil.rmtree(str(item))
            else:
                item.unlink()
            deleted += 1
        except Exception:
            pass
    return {"ok": True, "deleted": deleted}


@app.delete("/api/files/{filepath:path}")
async def delete_output_file(filepath: str):
    target = _safe_audio_path(filepath)
    if not target.exists():
        raise HTTPException(404, "Plik nie istnieje.")
    if target.is_dir():
        shutil.rmtree(str(target))
    else:
        target.unlink()
    return {"ok": True}


@app.get("/api/audio/{filepath:path}")
async def serve_audio(filepath: str):
    target = _safe_audio_path(filepath)
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="Plik nie istnieje.")
    return FileResponse(str(target))

@app.post("/api/translate-prompt")
async def translate_prompt(req: AIRequest):
    system_prompt = """Output ONLY valid JSON, no explanation. Use this exact structure: {"segments":[{"speaker_key":"pl_male_marek","text":"Hello","rate":"+0%","pitch":"+0Hz"},{"speaker_key":"pl_female_zofia","text":"Hi","rate":"-5%","pitch":"+2Hz"}],"silence_between_ms":500,"style":"normal"}. Available speaker_key values: pl_male_marek, pl_female_zofia, en_male_guy, en_male_andrew."""
    # Use a reliable non-thinking model for JSON generation
    _json_model = "qwen2.5:1.5b"
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                f"{req.ollama_host}/api/chat",
                json={
                    "model": _json_model,
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": req.prompt}
                    ],
                    "stream": False,
                }
            )
            import re as _re

            def _repair_json(s: str) -> str:
                """Fix common LLM JSON mistakes in rate/pitch/rate fields."""
                # "key"+N or "key"-N missing colon: "pitch"+2 -> "pitch":"+2"
                s = _re.sub(r'"(\w+)"\s*\+([\d.]+)', r'"\1":"+\2"', s)
                s = _re.sub(r'"(\w+)"\s*-([\d.]+)', r'"\1":"-\2"', s)
                # ":+N" bare positive number: "pitch":+2 -> "pitch":"+2"
                s = _re.sub(r'"(pitch|rate|volume)":\s*\+([\d.]+)', r'"\1":"+\2"', s)
                # ":N" bare integer for pitch/rate -> string
                s = _re.sub(r'"(pitch|rate|volume)":\s*(-?[\d.]+)(?=[,}\s])', r'"\1":"\2"', s)
                return s

            def _try_extract_json(raw: str) -> dict:
                cleaned = _re.sub(r'<think>.*?</think>', '', raw, flags=_re.DOTALL).strip()
                for text in (cleaned, raw):
                    s = text.find('{')
                    e = text.rfind('}')
                    if s == -1 or e == -1:
                        continue
                    snippet = text[s:e+1]
                    # Skip system prompt placeholders (e.g. [...])
                    if '"..."' in snippet or ': ...' in snippet or '[...]' in snippet:
                        continue
                    try:
                        return json.loads(snippet)
                    except json.JSONDecodeError:
                        # Try repairing common mistakes
                        try:
                            return json.loads(_repair_json(snippet))
                        except json.JSONDecodeError as je:
                            logger.warning("JSON repair failed: %s | snippet[:300]=%r", je, snippet[:300])
                            continue
                logger.warning("No JSON found in raw[:500]=%r", raw[:500])
                raise ValueError(f"Brak JSON w odpowiedzi modelu: {raw[:200]!r}")

            msg = response.json()["message"]
            # Ollama Qwen3: content has answer, thinking has reasoning (separate fields)
            raw_content = msg.get("content", "")
            raw_thinking = msg.get("thinking", "")
            # Try content first, then thinking field as fallback
            for raw in (raw_content, raw_thinking):
                if not raw:
                    continue
                try:
                    result = _try_extract_json(raw)
                    return {"success": True, "command": result}
                except (ValueError, json.JSONDecodeError):
                    continue
            raise ValueError(f"Brak JSON. content={raw_content[:100]!r} thinking={raw_thinking[:100]!r}")
    except httpx.TimeoutException:
        raise HTTPException(status_code=500, detail=f"Model {req.model} nie odpowiedział w czasie (timeout). Sprawdź czy Ollama działa i model jest załadowany.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e) or f"{type(e).__name__}: sprawdź logi serwera")

@app.get("/api/ollama/status")
async def check_ollama(host: str = "http://localhost:11434"):
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            r = await client.get(f"{host}/api/tags")
            if r.status_code == 200:
                models = [m["name"] for m in r.json().get("models", [])]
                return {"online": True, "models": models}
    except Exception:
        pass
    return {"online": False, "models": []}

@app.get("/api/ai/status")
async def check_ai_status():
    """Reports GPU and XTTS availability."""
    gpu = _check_gpu()
    return {
        "gpu_available": gpu,
        "xtts_loaded": _xtts_model is not None,
        "xtts_mode": "local",
    }


@app.get("/api/hardware")
async def get_hardware_info():
    """Returns full hardware profile + capability matrix for the UI."""
    import psutil as _psutil
    import torch as _torch

    gpu_avail = _torch.cuda.is_available()
    if gpu_avail:
        try:
            gpu_name = _torch.cuda.get_device_name(0)
            vram_total = _torch.cuda.get_device_properties(0).total_memory / 1e9
            vram_free_b, _ = _torch.cuda.mem_get_info(0)
            vram_free = vram_free_b / 1e9
        except Exception:
            gpu_name, vram_total, vram_free = "GPU", 0.0, 0.0
    else:
        gpu_name, vram_total, vram_free = "Brak", 0.0, 0.0

    ram = _psutil.virtual_memory()
    try:
        disk = _psutil.disk_usage(str(AUDIO_DIR))
    except Exception:
        disk = _psutil.disk_usage('/')
    cpu_count = _psutil.cpu_count(logical=False) or _psutil.cpu_count() or 1
    cpu_pct = _psutil.cpu_percent(interval=0.1)

    ram_free_gb = ram.available / 1e9

    cap = {
        "tts_fast": True,
        "tts_clone_xtts": gpu_avail or ram_free_gb >= 4,
        "video_generation": vram_total >= 6.0,
        "video_generation_hd": vram_total >= 10.0,
        "transcription": vram_free >= 3.0 or ram_free_gb >= 8.0,
        "music_generation": vram_total >= 5.0,
        "avatar_lipsync": vram_total >= 7.0,
    }

    def _time(cap_key: str, *, yes: str, warn: str, no: str) -> str:
        if cap.get(cap_key):
            return yes
        return warn if ram_free_gb >= 2 else no

    return {
        "gpu": {
            "name": gpu_name,
            "vram_total_gb": round(vram_total, 1),
            "vram_free_gb": round(vram_free, 1),
            "available": gpu_avail,
        },
        "ram": {
            "total_gb": round(ram.total / 1e9, 1),
            "free_gb": round(ram_free_gb, 1),
        },
        "cpu": {"cores": cpu_count, "usage_percent": round(cpu_pct, 1)},
        "disk": {"free_gb": round(disk.free / 1e9, 1)},
        "capabilities": cap,
        "estimated_times": {
            "tts_30s": "natychmiastowe",
            "clone_30s": "~30 sekund" if gpu_avail else "~90 sekund (CPU)",
            "video_5s": _time("video_generation", yes="~3 minuty", warn="~8 minut", no="niedostępne"),
            "audiobook_chapter": "~2 minuty" if gpu_avail else "~5 minut (CPU)",
            "transcription_1h": _time("transcription", yes="~3 minuty", warn="~10 minut", no="niedostępne"),
            "music_2min": _time("music_generation", yes="~2 minuty", warn="~8 minut", no="niedostępne"),
            "avatar_30s": _time("avatar_lipsync", yes="~4 minuty", warn="~15 minut", no="niedostępne"),
        },
        "installed_models": {
            "xtts_v2": True,
            "edge_tts": True,
            "animatediff": True,
            "melotts": _check_melotts(),
            "whisperx": _check_whisperx(),
            "musicgen": _check_musicgen(),
            "deepfilternet": _check_deepfilternet(),
            "demucs": _check_demucs(),
            "rvc": _check_rvc(),
            "echomimic": _check_echomimic(),
            "sadtalker": _check_sadtalker(),
        },
    }


@app.get("/api/gpu/queue/status")
async def get_gpu_queue_status():
    return GPU_QUEUE.status()


# ── Feature availability checks ──────────────────────────────────────────────
def _check_melotts() -> bool:
    try:
        import melo  # noqa: F401
        return True
    except ImportError:
        return False


def _check_whisperx() -> bool:
    try:
        import whisperx  # noqa: F401
        return True
    except ImportError:
        pass
    try:
        import whisper  # openai-whisper fallback
        return True
    except ImportError:
        return False


def _whisper_backend() -> str:
    """Returns 'whisperx' if available, 'whisper' if openai-whisper, else ''."""
    try:
        import whisperx  # noqa: F401
        return "whisperx"
    except ImportError:
        pass
    try:
        import whisper  # noqa: F401
        return "whisper"
    except ImportError:
        return ""


def _check_musicgen() -> bool:
    try:
        from audiocraft.models import MusicGen  # noqa: F401
        return True
    except ImportError:
        return False


def _check_deepfilternet() -> bool:
    try:
        import df  # noqa: F401
        return True
    except ImportError:
        return False


def _check_demucs() -> bool:
    try:
        import demucs  # noqa: F401
        return True
    except ImportError:
        return False


def _check_f5tts() -> bool:
    try:
        import f5_tts  # noqa: F401
        return True
    except ImportError:
        return False


def _check_cosyvoice() -> bool:
    try:
        import cosyvoice  # noqa: F401
        return True
    except ImportError:
        return False


# ── Transcription endpoint (WhisperX) ────────────────────────────────────────
@app.post("/api/transcribe")
async def transcribe_audio(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    language: str = Form("auto"),
    diarization: bool = Form(False),
):
    backend = _whisper_backend()
    if not backend:
        raise HTTPException(
            status_code=503,
            detail={
                "error": "whisperx_not_installed",
                "message": "WhisperX ani openai-whisper nie są zainstalowane.",
                "install_cmd": "pip install openai-whisper   # lub: pip install whisperx",
            },
        )
    suffix = Path(file.filename or "audio.mp3").suffix.lower() or ".mp3"
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    content = await file.read()
    tmp.write(content)
    tmp.close()

    job_id = uuid.uuid4().hex[:8]
    JOBS[job_id] = {"status": "processing", "progress": 5, "message": "Inicjalizacja transkrypcji…", "mode": "transcription"}

    async def _run():
        try:
            import torch as _t
            device = "cuda" if _t.cuda.is_available() else "cpu"

            if backend == "whisperx":
                import whisperx
                compute = "float16" if device == "cuda" else "int8"
                JOBS[job_id]["message"] = "Ładowanie modelu WhisperX large-v3…"
                model = whisperx.load_model("large-v3", device, compute_type=compute)
                audio = whisperx.load_audio(tmp.name)
                JOBS[job_id].update({"progress": 30, "message": "Transkrypcja (WhisperX)…"})
                result = model.transcribe(audio, batch_size=16, language=None if language == "auto" else language)
                JOBS[job_id].update({"progress": 70, "message": "Wyrównanie czasowe…"})
                model_a, metadata = whisperx.load_align_model(language_code=result["language"], device=device)
                result = whisperx.align(result["segments"], model_a, metadata, audio, device)
                if diarization:
                    JOBS[job_id].update({"progress": 85, "message": "Diaryzacja mówców…"})
                    try:
                        diarize_model = whisperx.DiarizationPipeline(use_auth_token=None, device=device)
                        diarize_segments = diarize_model(audio)
                        result = whisperx.assign_word_speakers(diarize_segments, result)
                    except Exception as e_d:
                        logger.warning("[Transcribe] Diarization failed: %s", e_d)
                segments = result.get("segments", [])
            else:
                # openai-whisper fallback (no diarization, no word-level timestamps)
                import whisper
                JOBS[job_id]["message"] = "Ładowanie modelu Whisper large-v3…"
                model = await asyncio.get_event_loop().run_in_executor(
                    None, lambda: whisper.load_model("large-v3", device=device)
                )
                JOBS[job_id].update({"progress": 30, "message": "Transkrypcja (openai-whisper)…"})
                lang_arg = None if language == "auto" else language
                result = await asyncio.get_event_loop().run_in_executor(
                    None, lambda: model.transcribe(tmp.name, language=lang_arg, verbose=False)
                )
                segments = [
                    {"start": s["start"], "end": s["end"], "text": s["text"]}
                    for s in result.get("segments", [])
                ]

            full_text = " ".join(s["text"] for s in segments)
            srt_lines = []
            for i, s in enumerate(segments, 1):
                def _ts(t):
                    h, r = divmod(int(t), 3600)
                    m, sec = divmod(r, 60)
                    ms = int((t - int(t)) * 1000)
                    return f"{h:02d}:{m:02d}:{sec:02d},{ms:03d}"
                speaker = f"[{s.get('speaker','?')}] " if "speaker" in s else ""
                srt_lines.append(f"{i}\n{_ts(s['start'])} --> {_ts(s['end'])}\n{speaker}{s['text'].strip()}\n")
            JOBS[job_id].update({
                "status": "completed", "progress": 100, "message": "Gotowe!",
                "text": full_text,
                "segments": segments,
                "srt": "\n".join(srt_lines),
                "language": result.get("language", language),
            })
        except Exception as e:
            JOBS[job_id].update({"status": "failed", "error": str(e)})
        finally:
            try:
                os.unlink(tmp.name)
            except OSError:
                pass

    background_tasks.add_task(_run)
    return {"job_id": job_id}


# ── Audio mixer endpoint (ffmpeg-based, always available) ────────────────────
@app.post("/api/audio/mix")
async def mix_audio(
    voice: UploadFile = File(...),
    music: Optional[UploadFile] = File(None),
    effects: Optional[UploadFile] = File(None),
    voice_volume: float = Form(1.0),
    music_volume: float = Form(0.3),
    effects_volume: float = Form(0.5),
    fade_in: bool = Form(False),
    fade_out: bool = Form(False),
):
    tmp_dir = Path(tempfile.mkdtemp(prefix="vs_mix_"))
    try:
        # Read all uploaded files
        voice_bytes = await voice.read()
        voice_ext = Path(voice.filename or "voice.mp3").suffix.lower() or ".mp3"
        voice_path = str(tmp_dir / f"voice{voice_ext}")
        Path(voice_path).write_bytes(voice_bytes)

        inputs = ["-i", voice_path]
        filter_parts = [f"[0:a]volume={voice_volume:.2f}[v0]"]
        mix_labels = ["[v0]"]
        idx = 1

        if music:
            music_bytes = await music.read()
            music_ext = Path(music.filename or "music.mp3").suffix.lower() or ".mp3"
            music_path = str(tmp_dir / f"music{music_ext}")
            Path(music_path).write_bytes(music_bytes)
            inputs += ["-i", music_path]
            filter_parts.append(f"[{idx}:a]volume={music_volume:.2f}[v{idx}]")
            mix_labels.append(f"[v{idx}]")
            idx += 1

        if effects:
            fx_bytes = await effects.read()
            fx_ext = Path(effects.filename or "fx.mp3").suffix.lower() or ".mp3"
            fx_path = str(tmp_dir / f"fx{fx_ext}")
            Path(fx_path).write_bytes(fx_bytes)
            inputs += ["-i", fx_path]
            filter_parts.append(f"[{idx}:a]volume={effects_volume:.2f}[v{idx}]")
            mix_labels.append(f"[v{idx}]")
            idx += 1

        n = len(mix_labels)
        amix_filter = f"{''.join(mix_labels)}amix=inputs={n}:duration=longest:normalize=0[mixed]"
        post = "[mixed]"
        post_filters = []
        if fade_in:
            post_filters.append("afade=t=in:d=2")
        if fade_out:
            post_filters.append("afade=t=out:d=3")
        final_label = "[out]"
        if post_filters:
            post_filter_str = ",".join(post_filters)
            amix_filter = amix_filter.replace("[mixed]", "") + "[mixed];" + f"[mixed]{post_filter_str}[out]"
        else:
            amix_filter = amix_filter.replace("[mixed]", "[out]")

        filter_complex = ";".join(filter_parts) + ";" + amix_filter

        out_name = f"mixed_{uuid.uuid4().hex[:8]}.mp3"
        out_path = str(AUDIO_DIR / out_name)

        cmd = ["ffmpeg", "-y"] + inputs + [
            "-filter_complex", filter_complex,
            "-map", "[out]", "-q:a", "4", out_path,
        ]
        r = subprocess.run(cmd, capture_output=True, timeout=120)
        if r.returncode != 0:
            raise HTTPException(500, detail=f"ffmpeg error: {r.stderr.decode()[-400:]}")

        return {"url": f"/api/audio/{out_name}", "filename": out_name}
    finally:
        shutil.rmtree(str(tmp_dir), ignore_errors=True)


# ── Noise reduction (DeepFilterNet) ─────────────────────────────────────────
@app.post("/api/audio/denoise")
async def denoise_audio(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not _check_deepfilternet():
        raise HTTPException(
            status_code=503,
            detail={
                "error": "deepfilternet_not_installed",
                "message": "DeepFilterNet nie jest zainstalowany.",
                "install_cmd": "pip install deepfilternet",
            },
        )
    suffix = Path(file.filename or "audio.wav").suffix.lower() or ".wav"
    tmp_in = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    tmp_in.write(await file.read())
    tmp_in.close()

    out_name = f"denoised_{uuid.uuid4().hex[:8]}.wav"
    out_path = str(AUDIO_DIR / out_name)
    job_id = uuid.uuid4().hex[:8]
    JOBS[job_id] = {"status": "processing", "progress": 10, "message": "Usuwanie szumów…", "mode": "denoise"}

    async def _run():
        try:
            import df  # noqa: F401
            from df.enhance import enhance, init_df, load_audio, save_audio
            model, df_state, _ = init_df()
            audio, _ = load_audio(tmp_in.name, sr=df_state.sr())
            enhanced = enhance(model, df_state, audio)
            save_audio(out_path, enhanced, df_state.sr())
            JOBS[job_id].update({
                "status": "completed", "progress": 100, "message": "Gotowe!",
                "url": f"/api/audio/{out_name}", "filename": out_name,
            })
        except Exception as e:
            JOBS[job_id].update({"status": "failed", "error": str(e)})
        finally:
            try:
                os.unlink(tmp_in.name)
            except OSError:
                pass

    background_tasks.add_task(_run)
    return {"job_id": job_id}


# ── Music generation (MusicGen) ──────────────────────────────────────────────
@app.post("/api/music/generate")
async def generate_music(background_tasks: BackgroundTasks, request: Request):
    body = await request.json()
    prompt = body.get("prompt", "background music")
    duration = min(max(int(body.get("duration", 30)), 5), 300)

    if not _check_musicgen():
        raise HTTPException(
            status_code=503,
            detail={
                "error": "musicgen_not_installed",
                "message": "MusicGen (audiocraft) nie jest zainstalowany.",
                "install_cmd": "pip install audiocraft",
            },
        )

    job_id = uuid.uuid4().hex[:8]
    JOBS[job_id] = {"status": "processing", "progress": 5, "message": "Ładowanie MusicGen…", "mode": "music"}

    async def _run():
        try:
            from audiocraft.models import MusicGen
            import torchaudio as _ta, torch as _t
            JOBS[job_id].update({"progress": 20, "message": "Generowanie muzyki…"})
            model = MusicGen.get_pretrained("facebook/musicgen-small")
            model.set_generation_params(duration=duration)
            wav = model.generate([prompt])
            out_name = f"music_{uuid.uuid4().hex[:8]}.wav"
            out_path = str(AUDIO_DIR / out_name)
            _ta.save(out_path, wav[0].cpu(), model.sample_rate)
            JOBS[job_id].update({
                "status": "completed", "progress": 100, "message": "Gotowe!",
                "url": f"/api/audio/{out_name}", "filename": out_name,
            })
        except Exception as e:
            JOBS[job_id].update({"status": "failed", "error": str(e)})

    background_tasks.add_task(_run)
    return {"job_id": job_id}

class VoiceSegmentSimple(BaseModel):
    voiceKey: str = "pl_male_marek"
    text: str = ""

class VideoRequest(BaseModel):
    prompt: str
    dialogue: str = ""
    speaker_key: str = "pl_male_marek"
    aspect_ratio: str = "256x256"
    with_audio: bool = False
    duration: int = 4
    video_style: str = "cinematic"
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"
    base_model: str = "emilianJR/epiCRealism"
    negative_prompt: str = "blurry, lowres, ugly, deformed"
    motion_lora: str = ""
    num_inference_steps: int = 25
    guidance_scale: float = 7.5
    voice_segments: list[VoiceSegmentSimple] = []
    engine: str = "animatediff"   # animatediff | wan21
    quality: str = "standard"     # fast | standard | high
    seed: int = -1                 # -1 = random

    @field_validator('aspect_ratio')
    @classmethod
    def validate_aspect_ratio(cls, v):
        parts = v.split('x')
        if len(parts) != 2: raise ValueError("Format: WxH")
        return v

@app.get("/api/jobs/{job_id}")
async def get_job_status(job_id: str):
    if job_id not in JOBS:
        raise HTTPException(status_code=404, detail="Job not found")
    return JOBS[job_id]

async def _process_wan21(job_id: str, req: VideoRequest, tmp_dir: Path) -> str:
    """Generate video using WAN 2.1 (Wan-AI/Wan2.1-T2V-1.3B — fits 12 GB VRAM)."""
    import torch
    try:
        from diffusers import WanPipeline
    except ImportError:
        raise Exception("WAN 2.1 wymaga diffusers >= 0.32. Uruchom: pip install -U diffusers")

    quality_steps = {"fast": 20, "standard": 40, "high": 60}
    steps = quality_steps.get(req.quality, 40)
    generator = torch.Generator("cuda").manual_seed(req.seed) if req.seed >= 0 else None

    JOBS[job_id]["message"] = "Ładowanie WAN 2.1 (Wan2.1-T2V-1.3B)…"
    torch.cuda.empty_cache()
    pipe = WanPipeline.from_pretrained(
        "Wan-AI/Wan2.1-T2V-1.3B",
        cache_dir=str(MODELS_DIR),
        torch_dtype=torch.bfloat16,
    ).to("cuda")

    JOBS[job_id]["message"] = "Generowanie wideo WAN 2.1…"
    JOBS[job_id]["progress"] = 30
    output = pipe(
        prompt=req.prompt,
        negative_prompt=req.negative_prompt,
        height=480, width=832,
        num_frames=81,
        num_inference_steps=steps,
        generator=generator,
    )
    out_path = tmp_dir / "wan21_output.mp4"
    output.frames[0][0].save(str(out_path))  # diffusers exports as video

    del pipe
    torch.cuda.empty_cache()

    final = AUDIO_DIR / f"video-wan21-{job_id}.mp4"
    shutil.copy(str(out_path), str(final))
    return f"/api/audio/{final.name}"


async def process_video_job(job_id: str, req: VideoRequest):
    JOBS[job_id] = {"status": "processing", "progress": 0, "message": "Planning storyboard..."}
    try:
        tmp_dir = AUDIO_DIR / f"vgen-{job_id}"
        tmp_dir.mkdir(parents=True, exist_ok=True)

        # WAN 2.1 shortcut — no storyboard needed
        if req.engine == "wan21":
            url = await _process_wan21(job_id, req, tmp_dir)
            JOBS[job_id] = {"status": "completed", "progress": 100, "url": url, "mode": "video_wan21"}
            return
        
        # 1. Planowanie
        num_chunks = max(1, req.duration // 2)
        storyboard = []
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                plan_prompt = f"Stwórz plan filmu ({num_chunks} scen) na podstawie opisu: {req.prompt}. Każda scena 2s. Styl: {req.video_style}. Zwróć JSON: {{\"scenes\": [{{ \"sd_prompt\": \"...\", \"text\": \"...\" }}]}}"
                res = await client.post(f"{req.ollama_host}/api/chat", json={
                    "model": req.model, "stream": False,
                    "messages": [{"role": "system", "content": "Expert director. JSON only."}, {"role": "user", "content": plan_prompt}]
                })
                storyboard = json.loads(res.json()["message"]["content"])["scenes"]
        except Exception as e:
            logger.warning("[process_video_job] Storyboard planning failed, using fallback: %s", e)
            storyboard = [{"sd_prompt": req.prompt, "text": req.dialogue}] * num_chunks

        style_suffix = STYLE_MAP.get(req.video_style, "")
        gpu = _check_gpu()

        if gpu:
            # ── GPU path: full AnimateDiff video generation ──────────────────
            clips = []
            import torch  # noqa: PLC0415
            from diffusers import AnimateDiffPipeline, MotionAdapter  # noqa: PLC0415
            from diffusers.utils import export_to_video  # noqa: PLC0415

            w, h = map(int, req.aspect_ratio.split('x'))
            scenes = storyboard[:num_chunks]

            # Phase 1: generate audio track
            # voice_segments = multi-speaker lines; fallback to single dialogue text
            JOBS[job_id]["message"] = "Generating audio tracks..."
            audio_map: dict = {}
            if req.voice_segments:
                # Concatenate all voice segments into a single audio file, assign to scene 0
                # (video will loop to match audio via ffmpeg -stream_loop)
                parts = []
                for vs in req.voice_segments:
                    if vs.text.strip():
                        try:
                            p = await generate_segment_audio(
                                SpeakerSegment(speaker_key=vs.voiceKey, text=vs.text), "normal", str(tmp_dir)
                            )
                            parts.append(p)
                        except Exception:
                            pass
                if parts:
                    combined = str(tmp_dir / "combined_audio.mp3")
                    list_f = str(tmp_dir / "audio_list.txt")
                    with open(list_f, "w") as f:
                        for p in parts:
                            f.write(f"file '{p}'\n")
                    proc = await asyncio.create_subprocess_exec(
                        "ffmpeg", "-f", "concat", "-safe", "0", "-i", list_f,
                        "-c:a", "libmp3lame", "-q:a", "2", combined, "-y",
                        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
                    )
                    await proc.communicate()
                    if os.path.exists(combined):
                        audio_map[0] = combined  # attach to first scene; ffmpeg will loop video
            elif req.with_audio:
                for i, scene in enumerate(scenes):
                    text = scene.get("text", "")
                    if text.strip():
                        try:
                            audio_map[i] = await generate_segment_audio(
                                SpeakerSegment(speaker_key=req.speaker_key, text=text), "normal", str(tmp_dir)
                            )
                        except Exception:
                            pass

            # Phase 2: unload Ollama model from VRAM, then load AnimateDiff
            os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"
            JOBS[job_id]["message"] = "Freeing GPU memory (unloading Ollama)..."
            try:
                async with httpx.AsyncClient(timeout=10.0) as _c:
                    await _c.post(f"{req.ollama_host}/api/generate",
                                  json={"model": req.model, "keep_alive": 0})
            except Exception:
                pass
            await asyncio.sleep(2)  # give Ollama time to unload
            torch.cuda.empty_cache()
            JOBS[job_id]["message"] = "Loading AnimateDiff pipeline..."
            adapter = MotionAdapter.from_pretrained(
                "guoyww/animatediff-motion-adapter-v1-5-2",
                cache_dir=str(MODELS_DIR), torch_dtype=torch.float16,
                low_cpu_mem_usage=True,
            )
            pipe = AnimateDiffPipeline.from_pretrained(
                req.base_model, motion_adapter=adapter,
                cache_dir=str(MODELS_DIR), torch_dtype=torch.float16,
                low_cpu_mem_usage=True,
            )
            if req.motion_lora:
                JOBS[job_id]["message"] = f"Loading motion LoRA ({req.motion_lora})..."
                pipe.load_lora_weights(
                    "guoyww/animatediff-motion-lora-" + req.motion_lora,
                    cache_dir=str(MODELS_DIR), adapter_name="motion",
                )
            pipe.enable_sequential_cpu_offload()

            # Phase 3: render each scene
            for i, scene in enumerate(scenes):
                JOBS[job_id]["progress"] = int((i / num_chunks) * 90)
                JOBS[job_id]["message"] = f"Rendering segment {i+1}/{num_chunks}..."
                try:
                    raw_vid = tmp_dir / f"raw_{i}.mp4"
                    output = pipe(
                        prompt=f"{scene['sd_prompt']}, {style_suffix}",
                        negative_prompt=req.negative_prompt,
                        num_frames=16, width=w, height=h,
                        num_inference_steps=req.num_inference_steps,
                        guidance_scale=req.guidance_scale,
                    )
                    export_to_video(output.frames[0], str(raw_vid))
                    chunk_final = tmp_dir / f"chunk_{i}.mp4"
                    temp_audio = audio_map.get(i, "")
                    if temp_audio and os.path.exists(temp_audio):
                        proc = await asyncio.create_subprocess_exec(
                            "ffmpeg", "-stream_loop", "-1", "-i", str(raw_vid),
                            "-i", str(temp_audio),
                            "-c:v", "libx264", "-c:a", "aac", "-shortest",
                            "-map", "0:v:0", "-map", "1:a:0",
                            str(chunk_final), "-y",
                            stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
                        )
                    else:
                        proc = await asyncio.create_subprocess_exec(
                            "ffmpeg", "-i", str(raw_vid), "-c:v", "libx264",
                            str(chunk_final), "-y",
                            stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
                        )
                    await proc.communicate()
                    if chunk_final.exists():
                        clips.append(str(chunk_final))
                except Exception as e:
                    logger.warning("[process_video_job] segment %d failed: %s", i, e)

            del pipe, adapter
            torch.cuda.empty_cache()

            if not clips:
                raise Exception("Żadna scena nie została wyrenderowana.")
            JOBS[job_id]["message"] = "Finalizing video..."
            final_filename = f"video-long-{job_id}.mp4"
            final_path = AUDIO_DIR / final_filename
            list_path = tmp_dir / "list.txt"
            with open(list_path, "w") as f:
                for c in clips:
                    f.write(f"file '{c}'\n")
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path), "-c", "copy", str(final_path), "-y",
            )
            await proc.communicate()
            if not final_path.exists():
                raise Exception("ffmpeg nie złączył klipów wideo.")
            JOBS[job_id] = {"status": "completed", "progress": 100, "url": f"/api/audio/{final_filename}", "mode": "video"}
        else:
            # ── CPU-only (no GPU): Audio Drama mode — combine scene audio into MP3 ──
            JOBS[job_id]["message"] = "GPU not available — generating Audio Drama (MP3)..."
            audio_parts = []
            for i, scene in enumerate(storyboard[:num_chunks]):
                JOBS[job_id]["progress"] = int((i / num_chunks) * 90)
                JOBS[job_id]["message"] = f"Generating audio {i+1}/{num_chunks}..."
                text = scene.get("text") or req.dialogue
                if text.strip():
                    part = await generate_segment_audio(
                        SpeakerSegment(speaker_key=req.speaker_key, text=text), "normal", str(tmp_dir),
                    )
                    audio_parts.append(part)

            if not audio_parts:
                raise Exception("Brak tekstu do wypowiedzenia — nie wygenerowano audio.")
            final_filename = f"video-audio-{job_id}.mp3"
            final_path = AUDIO_DIR / final_filename
            list_path = tmp_dir / "list.txt"
            with open(list_path, "w") as f:
                for p in audio_parts:
                    f.write(f"file '{p}'\n")
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path),
                "-c:a", "libmp3lame", "-q:a", "2", str(final_path), "-y",
            )
            await proc.communicate()
            if not final_path.exists():
                raise Exception("ffmpeg nie wygenerował pliku audio.")
            JOBS[job_id] = {
                "status": "completed", "progress": 100,
                "url": f"/api/audio/{final_filename}",
                "mode": "audio_drama",
                "warning": "GPU unavailable — audio-only output generated.",
            }
    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e)}
    finally:
        if tmp_dir.exists():
            shutil.rmtree(tmp_dir, ignore_errors=True)

@app.post("/api/generate-video")
async def generate_video(req: VideoRequest, background_tasks: BackgroundTasks):
    job_id = uuid.uuid4().hex[:8]
    background_tasks.add_task(process_video_job, job_id, req)
    return {"job_id": job_id}


async def process_movie_job(job_id: str, script: list, actors_meta: list, video_style: str, movie_dir: Path):
    JOBS[job_id] = {"status": "processing", "progress": 0, "message": "Initializing movie render..."}
    try:
        actor_paths = {}
        for actor in actors_meta:
            a_id = actor["id"]
            img_path = movie_dir / f"actor_{a_id}.png"
            if img_path.exists(): actor_paths[a_id] = img_path

        style_suffix = STYLE_MAP.get(video_style, "")
        gpu = _check_gpu()
        errors = []

        if gpu:
            # ── GPU path: full AnimateDiff ────────────────────────────────────
            clips = []
            import torch  # noqa: PLC0415
            from diffusers import AnimateDiffPipeline, MotionAdapter  # noqa: PLC0415
            from diffusers.utils import export_to_video  # noqa: PLC0415

            # Determine valid lines upfront
            valid_lines = [
                (i, line) for i, line in enumerate(script)
                if line.get("text", "").strip() and line.get("actorId") in actor_paths
            ]

            # Phase 1: generate all audio (XTTS always on CPU — no VRAM conflict)
            JOBS[job_id]["message"] = "Generating audio tracks..."
            audio_map: dict = {}
            for i, line in valid_lines:
                actor_info = next((a for a in actors_meta if a["id"] == line["actorId"]), {})
                voice_key = actor_info.get("voice_key", "pl_male_marek")
                try:
                    audio_map[i] = await generate_segment_audio(
                        SpeakerSegment(speaker_key=voice_key, text=line["text"]), "normal", str(movie_dir),
                    )
                except Exception as e:
                    errors.append(f"Audio line {i}: {e}")

            # Phase 2: unload Ollama model from VRAM, then load AnimateDiff
            os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"
            JOBS[job_id]["message"] = "Freeing GPU memory (unloading Ollama)..."
            try:
                async with httpx.AsyncClient(timeout=10.0) as _c:
                    await _c.post("http://localhost:11434/api/generate",
                                  json={"model": "qwen3.5:4b", "keep_alive": 0})
            except Exception:
                pass
            await asyncio.sleep(2)
            torch.cuda.empty_cache()
            JOBS[job_id]["message"] = "Loading AnimateDiff pipeline..."
            adapter = MotionAdapter.from_pretrained(
                "guoyww/animatediff-motion-adapter-v1-5-2",
                cache_dir=str(MODELS_DIR), torch_dtype=torch.float16,
                low_cpu_mem_usage=True,
            )
            pipe = AnimateDiffPipeline.from_pretrained(
                "emilianJR/epiCRealism", motion_adapter=adapter,
                cache_dir=str(MODELS_DIR), torch_dtype=torch.float16,
                low_cpu_mem_usage=True,
            )
            pipe.enable_sequential_cpu_offload()

            # Phase 3: render each valid line
            for step, (i, line) in enumerate(valid_lines):
                JOBS[job_id]["progress"] = int((step / len(valid_lines)) * 90)
                JOBS[job_id]["message"] = f"Rendering line {step+1}/{len(valid_lines)}..."
                actor_info = next((a for a in actors_meta if a["id"] == line["actorId"]), {})
                actor_name = actor_info.get("name", "Aktor")
                text = line["text"]
                try:
                    raw_vid = movie_dir / f"raw_{i}.mp4"
                    output = pipe(
                        prompt=f"{actor_name} speaking: {text[:100]}, {style_suffix}",
                        negative_prompt="lowres, blurry", num_frames=16, width=256, height=256,
                    )
                    export_to_video(output.frames[0], str(raw_vid))
                    clip_path = movie_dir / f"clip_{i}.mp4"
                    temp_audio = audio_map.get(i, "")
                    if temp_audio and os.path.exists(temp_audio):
                        # Loop video to match full audio length (audio drives clip duration)
                        proc = await asyncio.create_subprocess_exec(
                            "ffmpeg", "-stream_loop", "-1", "-i", str(raw_vid),
                            "-i", str(temp_audio),
                            "-c:v", "libx264", "-c:a", "aac", "-shortest",
                            "-map", "0:v:0", "-map", "1:a:0",
                            str(clip_path), "-y",
                            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                        )
                    else:
                        proc = await asyncio.create_subprocess_exec(
                            "ffmpeg", "-i", str(raw_vid), "-c:v", "libx264", "-t", "2", str(clip_path), "-y",
                            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                        )
                    _, ffmpeg_err = await proc.communicate()
                    if clip_path.exists():
                        clips.append(str(clip_path))
                    else:
                        err_msg = ffmpeg_err.decode()[-200:] if ffmpeg_err else "unknown"
                        errors.append(f"ffmpeg failed line {i}: {err_msg}")
                        logger.warning("[movie] ffmpeg failed line %d: %s", i, err_msg)
                except Exception as e:
                    logger.exception("[movie] Exception line %d: %s", i, e)
                    errors.append(f"Line {i}: {str(e)}")

            del pipe, adapter
            torch.cuda.empty_cache()

            if not clips:
                raise Exception(f"No clips generated. Errors: {errors[:3]}")

            JOBS[job_id]["message"] = "Finalizing movie..."
            final_filename = f"movie-{job_id}.mp4"
            final_path = AUDIO_DIR / final_filename
            list_path = movie_dir / "clips.txt"
            with open(list_path, "w") as f:
                for c in clips:
                    f.write(f"file '{c}'\n")
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path), "-c", "copy", str(final_path), "-y",
            )
            await proc.communicate()
            JOBS[job_id] = {
                "status": "completed", "progress": 100,
                "url": f"/api/audio/{final_filename}",
                "mode": "video", "warnings": errors,
            }
        else:
            # ── CPU-only: Audio Drama — one MP3 per actor line, then concat ──
            JOBS[job_id]["message"] = "GPU not available — generating Audio Drama (MP3)..."
            audio_parts = []
            for i, line in enumerate(script):
                JOBS[job_id]["progress"] = int((i / len(script)) * 90)
                JOBS[job_id]["message"] = f"Generating audio {i+1}/{len(script)}..."
                a_id, text = line["actorId"], line["text"]
                if not text.strip():
                    continue
                actor_info = next((a for a in actors_meta if a["id"] == a_id), {})
                voice_key = actor_info.get("voice_key", "pl_male_marek")
                try:
                    part = await generate_segment_audio(
                        SpeakerSegment(speaker_key=voice_key, text=text), "dramatic", str(movie_dir),
                    )
                    audio_parts.append(part)
                except Exception as e:
                    errors.append(f"Line {i}: {str(e)}")

            final_filename = f"movie-audio-{job_id}.mp3"
            final_path = AUDIO_DIR / final_filename
            if audio_parts:
                list_path = movie_dir / "clips.txt"
                with open(list_path, "w") as f:
                    for p in audio_parts:
                        f.write(f"file '{p}'\n")
                proc = await asyncio.create_subprocess_exec(
                    "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path),
                    "-c:a", "libmp3lame", "-q:a", "2", str(final_path), "-y",
                )
                await proc.communicate()
            else:
                raise Exception("No audio lines generated")

            JOBS[job_id] = {
                "status": "completed", "progress": 100,
                "url": f"/api/audio/{final_filename}",
                "mode": "audio_drama",
                "warning": "GPU unavailable — audio-only output generated.",
                "warnings": errors,
            }
    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e)}
    finally:
        if movie_dir.exists():
            shutil.rmtree(movie_dir, ignore_errors=True)

@app.post("/api/render-movie")
async def render_movie(request: Request, background_tasks: BackgroundTasks):
    form = await request.form()
    script = json.loads(form.get("script", "[]"))
    actors_meta = json.loads(form.get("actors_meta", "[]"))
    video_style = form.get("video_style", "cinematic")
    try:
        max_lines = max(1, int(form.get("duration", 60)) // 3)
    except (ValueError, TypeError):
        max_lines = 20
    job_id = uuid.uuid4().hex[:8]
    movie_dir = AUDIO_DIR / f"movie-{job_id}"
    movie_dir.mkdir(parents=True, exist_ok=True)

    # Save images before background task
    for actor in actors_meta:
        a_id = actor["id"]
        if f"actor_image_{a_id}" in form:
            img_path = movie_dir / f"actor_{a_id}.png"
            with open(img_path, "wb") as f: f.write(await form[f"actor_image_{a_id}"].read())

    background_tasks.add_task(process_movie_job, job_id, script[:max_lines], actors_meta, video_style, movie_dir)
    return {"job_id": job_id}


@app.get("/api/health")
async def health():
    return {"status": "ok", "version": "2.0.0"}

_ALLOWED_AUDIO_TYPES = {"audio/wav", "audio/wave", "audio/mpeg", "audio/mp3", "audio/x-wav"}
_MAX_SPEAKER_SIZE = 10 * 1024 * 1024  # 10 MB

@app.post("/api/clone-voice")
async def clone_voice(text: str = Form(...), speaker_name: str = Form(...), file: Optional[UploadFile] = File(None), voice_lang: str = Form('pl')):
    if file:
        if file.content_type and file.content_type not in _ALLOWED_AUDIO_TYPES:
            raise HTTPException(status_code=415, detail=f"Niedozwolony typ pliku: {file.content_type}. Wymagane: WAV/MP3")
        audio_bytes = await file.read()
        if len(audio_bytes) > _MAX_SPEAKER_SIZE:
            raise HTTPException(status_code=413, detail="Plik za duży (max 10 MB)")
    else:
        audio_bytes = None

    safe_name = "".join([c for c in speaker_name if c.isalnum() or c in ("-", "_")]) or f"spk_{uuid.uuid4().hex[:6]}"
    safe_name = safe_name[:50]

    existing_voice = not bool(audio_bytes)

    if existing_voice:
        # No file uploaded — look up already-saved voice
        speaker_path = None
        for ext in (".wav", ".mp3"):
            candidate = XTTS_SPEAKERS_DIR / (safe_name + ext)
            if candidate.exists():
                speaker_path = candidate
                break
        if not speaker_path:
            raise HTTPException(status_code=400, detail="Nie przesłano pliku próbki głosu.")
        # Return immediately — voice already exists, no preview needed
        stored_lang = _load_voice_lang(safe_name) or voice_lang
        return {
            "saved": True,
            "voice_name": safe_name,
            "filename": speaker_path.name,
            "voice_key": f"cloned_{safe_name}",
            "preview_url": None,
            "xtts_available": True,
            "lang": stored_lang,
        }

    speaker_path = XTTS_SPEAKERS_DIR / f"{safe_name}.wav"
    # Save raw upload to temp file, then convert to 22050Hz PCM WAV for XTTS
    import tempfile as _tf
    raw_suffix = ".wav"
    with _tf.NamedTemporaryFile(suffix=raw_suffix, delete=False) as tmp:
        tmp.write(audio_bytes)
        tmp_path = tmp.name
    try:
        _convert_to_xtts_wav(tmp_path, str(speaker_path))
    except Exception as conv_err:
        logger.warning("[clone_voice] ffmpeg conversion failed (%s), saving raw", conv_err)
        speaker_path.write_bytes(audio_bytes)
    finally:
        try:
            import os as _os
            _os.unlink(tmp_path)
        except OSError:
            pass

    # Save language metadata
    _save_voice_meta(safe_name, voice_lang)

    # Generate XTTS preview for newly uploaded voice
    preview_url: Optional[str] = None
    if text.strip():
        try:
            job_id = uuid.uuid4().hex[:8]
            wav_tmp = AUDIO_DIR / f"clone-{job_id}.wav"
            output_filename = f"clone-{job_id}.mp3"
            output_path = AUDIO_DIR / output_filename
            await xtts_synthesize(text, str(speaker_path), str(wav_tmp), force_language=voice_lang)
            _ffmpeg_wav_to_mp3(str(wav_tmp), str(output_path))
            if wav_tmp.exists():
                wav_tmp.unlink()
            preview_url = f"/api/audio/{output_filename}"
        except Exception as e:
            logger.warning("[clone_voice] XTTS preview failed: %s", e)

    return {
        "saved": True,
        "voice_name": safe_name,
        "filename": speaker_path.name,
        "voice_key": f"cloned_{safe_name}",
        "preview_url": preview_url,
        "xtts_available": preview_url is not None,
        "lang": voice_lang,
    }

@app.get("/api/cloned-voices")
async def list_cloned_voices():
    voices = []
    if XTTS_SPEAKERS_DIR.exists():
        for f in XTTS_SPEAKERS_DIR.glob("*"):
            if f.suffix.lower() in (".wav", ".mp3"):
                lang = _load_voice_lang(f.stem) or 'pl'
                voices.append({"name": f.stem, "filename": f.name, "lang": lang})
    return {"voices": voices}

@app.delete("/api/cloned-voices/{voice_name}")
async def delete_cloned_voice(voice_name: str):
    safe_name = Path(voice_name).name
    for ext in (".wav", ".mp3"):
        voice_path = XTTS_SPEAKERS_DIR / (safe_name + ext)
        if voice_path.exists():
            voice_path.unlink()
            meta_path = XTTS_SPEAKERS_DIR / f"{safe_name}.json"
            if meta_path.exists():
                meta_path.unlink()
            return {"status": "ok", "deleted": voice_name}
    raise HTTPException(status_code=404, detail=f"Głos '{voice_name}' nie znaleziony")

@app.post("/api/movie/generate-script")
async def generate_movie_script(req: AIRequest):
    system_prompt = """You are a professional movie script writer.
Create a script with multiple actors based on user prompt.
Output ONLY JSON:
{
  "actors": ["Actor Name 1", "Actor Name 2"],
  "script": [
    {"actorId": 1, "text": "..."},
    {"actorId": 2, "text": "..."}
  ]
}
Actor IDs must correspond to the index in 'actors' list (1-based).
"""
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(f"{req.ollama_host}/api/chat", json={
                "model": req.model, "stream": False,
                "messages": [{"role": "system", "content": system_prompt},
                             {"role": "user", "content": req.prompt}]
            })
            return json.loads(res.json()["message"]["content"])
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/video/generate-params")
async def generate_video_params(req: AIRequest):
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(f"{req.ollama_host}/api/chat", json={
                "model": req.model,
                "messages": [{"role": "system", "content": "Expert video AI. Output JSON: visual_prompt, dialogue"},
                             {"role": "user", "content": req.prompt}]
            })
            return json.loads(res.json()["message"]["content"])
    except Exception:
        return {"visual_prompt": req.prompt, "dialogue": ""}

_ECHOMIMIC_DIR = Path.home() / "EchoMimic"
_ECHOMIMIC_VENV_PY = (_EXTERNAL_ROOT / "echomimic-env" / "bin" / "python3") if _EXTERNAL_ROOT else Path("/opt/echomimic-env/bin/python3")

_SADTALKER_DIR = Path.home() / "SadTalker"
_SADTALKER_VENV_PY = (_EXTERNAL_ROOT / "sadtalker-env" / "bin" / "python3") if _EXTERNAL_ROOT else Path("/opt/sadtalker-env/bin/python3")


def _get_echomimic_python() -> str:
    if _ECHOMIMIC_VENV_PY.exists():
        return str(_ECHOMIMIC_VENV_PY)
    return "python3"


def _check_echomimic() -> bool:
    return (_ECHOMIMIC_DIR / "infer_audio2vid.py").exists()


async def _run_echomimic_clip(photo_path: Path, audio_path: Path, result_dir: Path) -> Path:
    """Animate photo with EchoMimic (diffusion-based lip sync). Returns h264 MP4."""
    result_dir.mkdir(parents=True, exist_ok=True)

    # EchoMimic requires 16 kHz mono WAV
    wav_path = result_dir / "audio.wav"
    conv = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", str(audio_path),
        "-ar", "16000", "-ac", "1", str(wav_path),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await conv.communicate()
    if not wav_path.exists():
        wav_path = audio_path

    # EchoMimic reads image/audio from YAML test_cases — generate per-job config
    weights_dir = _ECHOMIMIC_DIR / "pretrained_weights"
    config_yaml = f"""pretrained_base_model_path: "{weights_dir}/sd-image-variations-diffusers/"
pretrained_vae_path: "{weights_dir}/sd-vae-ft-mse/"
audio_model_path: "{weights_dir}/audio_processor/whisper_tiny.pt"
denoising_unet_path: "{weights_dir}/denoising_unet.pth"
reference_unet_path: "{weights_dir}/reference_unet.pth"
face_locator_path: "{weights_dir}/face_locator.pth"
motion_module_path: "{weights_dir}/motion_module.pth"
inference_config: "{_ECHOMIMIC_DIR}/configs/inference/inference_v2.yaml"
weight_dtype: 'fp16'
test_cases:
  "{photo_path.resolve()}":
    - "{wav_path.resolve()}"
"""
    config_path = result_dir / "em_config.yaml"
    config_path.write_text(config_yaml)

    env = os.environ.copy()
    env["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True,max_split_size_mb:256"
    # Enable CPU offload when free VRAM < 9 GB (other models like XTTS may be loaded)
    try:
        import subprocess as _sp
        vram_free_mb = int(_sp.check_output(
            ["nvidia-smi", "--query-gpu=memory.free", "--format=csv,noheader,nounits"],
            text=True, timeout=5,
        ).strip().split("\n")[0])
        if vram_free_mb < 9_000:
            env["ECHOMIMIC_CPU_OFFLOAD"] = "1"
    except Exception:
        env["ECHOMIMIC_CPU_OFFLOAD"] = "1"  # safe default

    python_bin = _get_echomimic_python()
    # CWD = result_dir so relative "output/..." goes inside result_dir
    proc = await asyncio.create_subprocess_exec(
        python_bin, str(_ECHOMIMIC_DIR / "infer_audio2vid.py"),
        "--config", str(config_path),
        "-W", "512", "-H", "512",
        "--seed", "42",
        "--steps", "30",
        "--fps", "25",
        env=env, cwd=str(result_dir),
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        _, stderr = await asyncio.wait_for(proc.communicate(), timeout=1200)
    except asyncio.TimeoutError:
        proc.kill()
        raise Exception("EchoMimic timeout (>20 min)")
    if proc.returncode != 0:
        raise Exception(f"EchoMimic error: {stderr.decode(errors='replace')[-500:]}")

    # Output lands in result_dir/output/**/..._withaudio.mp4
    candidates = sorted(
        result_dir.glob("**/*withaudio.mp4"),
        key=lambda p: p.stat().st_mtime, reverse=True,
    )
    if not candidates:
        # Fallback: any mp4 generated
        candidates = sorted(result_dir.glob("**/*.mp4"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not candidates:
        raise Exception("EchoMimic nie wygenerował MP4")
    raw_video = candidates[0]

    # Transcode to h264 + faststart for browser playback
    h264_video = result_dir / "echomimic_h264.mp4"
    proc2 = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", str(raw_video),
        "-c:v", "libx264", "-crf", "18", "-preset", "medium", "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        "-c:a", "aac", "-b:a", "192k",
        str(h264_video),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await proc2.communicate()
    return h264_video if h264_video.exists() else raw_video


def _get_sadtalker_python() -> str:
    """Return the Python executable to use for SadTalker (dedicated venv preferred)."""
    if _SADTALKER_VENV_PY.exists():
        return str(_SADTALKER_VENV_PY)
    return "python3"


def _check_sadtalker() -> bool:
    """Returns True if SadTalker inference script is present (venv optional)."""
    return (_SADTALKER_DIR / "inference.py").exists()


async def _run_sadtalker_clip(photo_path: Path, audio_path: Path, result_dir: Path) -> Path:
    """Audio-driven realistic talking head: SadTalker lip-sync + natural head motion +
    GFPGAN face restoration. Returns path to h264 MP4 with audio muxed."""
    result_dir.mkdir(parents=True, exist_ok=True)

    # SadTalker works best with 16 kHz mono WAV
    wav_path = result_dir / "audio.wav"
    conv = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", str(audio_path),
        "-ar", "16000", "-ac", "1", str(wav_path),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await conv.communicate()
    if not wav_path.exists():
        wav_path = audio_path  # fallback: pass original

    env = os.environ.copy()
    env["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"

    python_bin = _get_sadtalker_python()
    proc = await asyncio.create_subprocess_exec(
        python_bin, "inference.py",
        "--driven_audio", str(wav_path),
        "--source_image", str(photo_path),
        "--result_dir", str(result_dir),
        "--enhancer", "RestoreFormer",  # preserves more natural skin detail than GFPGAN
        "--preprocess", "crop",         # faster, better for headshots
        "--size", "256",                # 256 uses ~3 GB VRAM vs 512 which OOMs when XTTS loaded
        # no --still: natural head nods/motion avoids the "mannequin" look
        "--expression_scale", "1.1",
        env=env, cwd=str(_SADTALKER_DIR),
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        _, stderr = await asyncio.wait_for(proc.communicate(), timeout=900)
    except asyncio.TimeoutError:
        proc.kill()
        raise Exception("SadTalker timeout (>15 min) — tekst zbyt długi dla jednego klipu avatara")
    if proc.returncode != 0:
        raise Exception(f"SadTalker error: {stderr.decode(errors='replace')[-400:]}")

    candidates = sorted(result_dir.glob("*.mp4"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not candidates:
        raise Exception("SadTalker nie wygenerował MP4")
    raw_video = candidates[0]

    # SadTalker outputs mpeg4 which browsers can't play — transcode to h264
    final_clip = result_dir / "final_h264.mp4"
    proc2 = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", str(raw_video),
        "-c:v", "libx264", "-crf", "18", "-preset", "medium", "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        "-c:a", "aac", "-b:a", "192k",
        str(final_clip),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await proc2.communicate()
    return final_clip if final_clip.exists() else raw_video


async def _make_static_talking_clip(photo_path: Path, audio_path: Path, result_dir: Path) -> Path:
    """Fallback: pan-zoom photo + audio → MP4 (no face animation, no GPU needed)."""
    result_dir.mkdir(parents=True, exist_ok=True)
    out = result_dir / "talking.mp4"
    # Ken-Burns: slow zoom in + audio
    proc = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y",
        "-loop", "1", "-framerate", "25", "-i", str(photo_path),
        "-i", str(audio_path),
        "-vf", "scale=720:720:force_original_aspect_ratio=decrease,pad=720:720:(ow-iw)/2:(oh-ih)/2,zoompan=z='min(zoom+0.001,1.3)':d=1:x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)',format=yuv420p",
        "-c:v", "libx264", "-c:a", "aac", "-shortest",
        "-map", "0:v:0", "-map", "1:a:0",
        str(out),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await proc.communicate()
    if not out.exists():
        raise Exception("ffmpeg nie stworzył klipu fallback")
    return out


async def process_sadtalker_job(
    job_id: str,
    photo1: Path,
    photo2: Path,
    lines: list,          # [{"actor": 1|2, "text": str, "voice_id": str}]
    out_dir: Path,
):
    use_em = _check_echomimic()
    use_st = _check_sadtalker()
    if use_em:
        method = "EchoMimic"
    elif use_st:
        method = "SadTalker"
    else:
        method = "animacja statyczna"
    JOBS[job_id] = {"status": "processing", "progress": 0, "message": f"Przygotowanie ({method})…"}
    clips = []
    try:
        total = len(lines)

        for i, line in enumerate(lines):
            JOBS[job_id]["progress"] = int(i / total * 85)
            JOBS[job_id]["message"] = f"Kwestia {i+1}/{total}: generowanie audio…"

            audio_path = out_dir / f"line_{i}.mp3"
            communicate = edge_tts.Communicate(
                text=line["text"], voice=line["voice_id"], rate="+0%", pitch="+0Hz", volume="+0%"
            )
            await communicate.save(str(audio_path))

            JOBS[job_id]["message"] = f"Kwestia {i+1}/{total}: animacja ({method})…"
            photo = photo1 if line["actor"] == 1 else photo2
            clip_result_dir = out_dir / f"clip_{i}"
            clip_path = None
            try:
                if use_em:
                    clip_path = await _run_echomimic_clip(photo, audio_path, clip_result_dir)
                elif use_st:
                    clip_path = await _run_sadtalker_clip(photo, audio_path, clip_result_dir)
                else:
                    clip_path = await _make_static_talking_clip(photo, audio_path, clip_result_dir)
                clips.append(str(clip_path))
            except Exception as e:
                logger.warning("[Avatar] kwestia %d błąd (%s): %s", i, method, e)
                # Fallback: static animation
                try:
                    clip_path = await _make_static_talking_clip(photo, audio_path, clip_result_dir)
                    clips.append(str(clip_path))
                except Exception as e2:
                    logger.warning("[Avatar] fallback też zawiódł: %s", e2)

        if not clips:
            raise Exception("Żaden klip nie został wygenerowany.")

        JOBS[job_id]["message"] = "Scalanie klipów..."
        final_filename = f"dialogue-{job_id}.mp4"
        final_path = AUDIO_DIR / final_filename
        list_path = out_dir / "list.txt"
        with open(list_path, "w") as f:
            for c in clips:
                f.write(f"file '{c}'\n")
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path),
            "-c", "copy", str(final_path), "-y",
            stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
        )
        await proc.communicate()
        if not final_path.exists():
            raise Exception("ffmpeg nie złączył klipów.")

        JOBS[job_id] = {
            "status": "completed", "progress": 100,
            "url": f"/api/audio/{final_filename}",
            "filename": final_filename,
            "mode": "dialogue",
        }
    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e)}
    finally:
        if out_dir.exists():
            shutil.rmtree(out_dir, ignore_errors=True)


class DialogueScriptRequest(BaseModel):
    prompt: str
    actor1_name: str = "Aktor 1"
    actor2_name: str = "Aktor 2"
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"


def _extract_json(text: str) -> dict:
    """Robustly extract JSON from Ollama response (handles markdown fences, trailing commas)."""
    import re
    # Strip markdown code fences
    text = re.sub(r"^```(?:json)?\s*", "", text.strip(), flags=re.IGNORECASE)
    text = re.sub(r"\s*```$", "", text.strip())
    # Remove JS-style comments
    text = re.sub(r"//[^\n]*", "", text)
    text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
    # Remove trailing commas before } or ]
    text = re.sub(r",\s*([}\]])", r"\1", text)
    return json.loads(text)


@app.post("/api/dialogue/generate-script")
async def generate_dialogue_script(req: DialogueScriptRequest):
    system_prompt = (
        f"Jestes scenarzystem. Napisz dialog miedzy dwiema osobami.\n"
        f"Aktor 1: {req.actor1_name}\n"
        f"Aktor 2: {req.actor2_name}\n"
        f'Zwroc TYLKO JSON: {{"lines": [{{"actor": 1, "text": "..."}}, {{"actor": 2, "text": "..."}}]}}\n'
        f"Dialog: 4-8 kwestii, po polsku, naturalnie. Bez komentarzy, bez markdownu."
    )
    import traceback as _tb
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            try:
                res = await client.post(f"{req.ollama_host}/api/chat", json={
                    "model": req.model, "stream": False,
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": req.prompt},
                    ]
                })
                res.raise_for_status()
            except Exception as conn_err:
                err_detail = f"Błąd połączenia z Ollama ({req.ollama_host}): {conn_err}\n{_tb.format_exc()[-600:]}"
                logger.warning("[dialogue] %s", err_detail)
                raise HTTPException(status_code=503, detail=err_detail)
            try:
                raw = res.json()["message"]["content"]
            except Exception as parse_err:
                err_detail = f"Nieprawidłowa odpowiedź Ollama: {res.text[:300]} | {parse_err}"
                logger.warning("[dialogue] %s", err_detail)
                raise HTTPException(status_code=500, detail=err_detail)
            try:
                data = _extract_json(raw)
            except Exception:
                m = re.search(r'\{.*\}', raw, re.DOTALL)
                if m:
                    data = _extract_json(m.group(0))
                else:
                    raise HTTPException(status_code=500, detail=f"Model zwrócił niepoprawny JSON: {raw[:200]}")
            if "lines" not in data:
                raise HTTPException(status_code=500, detail=f"Brak klucza 'lines' w odpowiedzi: {raw[:200]}")
            return data
    except HTTPException:
        raise
    except Exception as e:
        err_detail = f"{type(e).__name__}: {e}\n{_tb.format_exc()[-600:]}"
        logger.error("[dialogue] Nieoczekiwany błąd: %s", err_detail)
        raise HTTPException(status_code=500, detail=err_detail)


class AnimationPromptRequest(BaseModel):
    description: str
    style: str = "anime"
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"


@app.post("/api/animation/generate-prompt")
async def generate_animation_prompt(req: AnimationPromptRequest):
    system_prompt = (
        f"You are an expert Stable Diffusion prompt engineer specializing in {req.style} animation.\n"
        f"Given a scene description, generate an optimized SD prompt and negative prompt.\n"
        f"Output ONLY JSON: {{\"prompt\": \"...\", \"negative_prompt\": \"...\"}}\n"
        f"The prompt should be detailed, comma-separated tags, English, max 120 words.\n"
        f"The negative_prompt should list things to avoid.\n"
        f"No markdown, no explanations, just JSON."
    )
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(f"{req.ollama_host}/api/chat", json={
                "model": req.model, "stream": False,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": req.description},
                ]
            })
            raw = res.json()["message"]["content"]
            try:
                data = _extract_json(raw)
            except Exception:
                import re
                m = re.search(r'\{.*\}', raw, re.DOTALL)
                data = _extract_json(m.group(0)) if m else {}
            if "prompt" not in data:
                raise HTTPException(status_code=500, detail=f"Brak 'prompt' w odpowiedzi: {raw[:200]}")
            return data
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


async def process_avatar_job(
    job_id: str,
    photo_path: Path,
    text: str,
    voice_key: str,
    out_dir: Path,
    preferred_model: str = "auto",
) -> None:
    """Standalone avatar: photo + text → TTS → EchoMimic/SadTalker/static → MP4.

    preferred_model: 'auto' | 'echomimic' | 'sadtalker' | 'wav2lip' | 'liveportrait' | 'musetalk'
    """
    use_em = _check_echomimic()
    use_st = _check_sadtalker()

    # Respect user's model preference; fallback chain if preferred model not installed
    if preferred_model == "echomimic":
        use_em, use_st = use_em, False         # force echomimic if installed, else static
    elif preferred_model == "sadtalker":
        use_em, use_st = False, use_st         # force sadtalker if installed, else static
    elif preferred_model in ("wav2lip", "liveportrait", "musetalk"):
        # These models are in the catalog but not yet integrated — fall through to auto
        logger.info("[Avatar] Model '%s' nie jest jeszcze zintegrowany, używam auto-doboru", preferred_model)

    if use_em:
        method = "EchoMimic"
    elif use_st:
        method = "SadTalker"
    else:
        method = "animacja statyczna"
    JOBS[job_id] = {"status": "processing", "progress": 5, "message": "Generowanie audio…"}

    try:
        catalog = await get_full_catalog()
        voice_info = catalog.get(voice_key)

        audio_path = out_dir / "avatar_audio.mp3"

        if voice_info and voice_info.get("type") == "cloned":
            ref = str(XTTS_SPEAKERS_DIR / voice_info["id"])
            voice_lang = voice_info.get("voice_lang") or _load_voice_lang(
                voice_info["id"].rsplit(".", 1)[0]
            ) or "pl"
            try:
                await chatterbox_synthesize(text, ref, str(audio_path), language=voice_lang)
            except Exception as _cb_err:
                logger.warning("[Avatar] Chatterbox fallback XTTS: %s", _cb_err)
                await xtts_synthesize(text, ref, str(audio_path), force_language=voice_lang)
        else:
            vid = (voice_info or {}).get("id", "pl-PL-MarekNeural")
            await edge_tts.Communicate(text, vid).save(str(audio_path))

        JOBS[job_id].update({"progress": 40, "message": f"Animacja twarzy ({method})…"})

        # Free XTTS from VRAM before spawning GPU-heavy face-animation subprocess.
        # This recovers ~3-4 GB and lets EchoMimic/SadTalker run on GPU without offload.
        _release_tts_vram()

        final_clip = None
        if use_em:
            try:
                final_clip = await _run_echomimic_clip(photo_path, audio_path, out_dir / "em")
            except Exception as e:
                logger.warning("[Avatar] EchoMimic error: %s — trying SadTalker", e)
        if final_clip is None and use_st:
            try:
                final_clip = await _run_sadtalker_clip(photo_path, audio_path, out_dir / "st")
            except Exception as e:
                logger.warning("[Avatar] SadTalker error: %s — trying static fallback", e)
        if final_clip is None:
            final_clip = await _make_static_talking_clip(photo_path, audio_path, out_dir / "static")

        final_filename = f"avatar-{job_id}.mp4"
        final_path = AUDIO_DIR / final_filename
        import shutil as _shutil
        _shutil.copy2(str(final_clip), str(final_path))

        JOBS[job_id] = {
            "status": "completed", "progress": 100,
            "url": f"/api/audio/{final_filename}",
            "filename": final_filename,
            "mode": "avatar",
        }
    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e)}
    finally:
        if out_dir.exists():
            import shutil as _shutil
            _shutil.rmtree(out_dir, ignore_errors=True)


@app.post("/api/render-avatar")
async def render_avatar(
    request: Request,
    background_tasks: BackgroundTasks,
):
    form = await request.form()

    photo: UploadFile = form.get("photo")
    if not photo:
        raise HTTPException(status_code=400, detail="Wymagane zdjęcie (photo).")
    allowed_img = {"image/jpeg", "image/png", "image/webp"}
    if photo.content_type and photo.content_type not in allowed_img:
        raise HTTPException(status_code=415, detail="Wymagane zdjęcie JPG/PNG/WEBP.")

    text = (form.get("text") or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="Pole text jest wymagane.")

    voice_key = form.get("voice_key", "pl_male_marek")
    avatar_model = form.get("avatar_model", "auto")

    job_id = uuid.uuid4().hex[:8]
    work_dir = AUDIO_DIR / f"avatar-{job_id}"
    work_dir.mkdir(parents=True, exist_ok=True)

    photo_suffix = Path(photo.filename or "photo.jpg").suffix or ".jpg"
    photo_path = work_dir / f"photo{photo_suffix}"
    photo_path.write_bytes(await photo.read())

    JOBS[job_id] = {"status": "queued", "progress": 0, "message": "Oczekuje w kolejce…", "mode": "avatar"}
    background_tasks.add_task(_enqueue_job, job_id, process_avatar_job, job_id, photo_path, text, voice_key, work_dir, avatar_model)
    return {"job_id": job_id}


@app.post("/api/render-dialogue")
async def render_dialogue(
    request: Request,
    background_tasks: BackgroundTasks,
):
    form = await request.form()

    allowed_img = {"image/jpeg", "image/png", "image/webp"}
    actor1_photo: UploadFile = form.get("actor1_photo")
    actor2_photo: UploadFile = form.get("actor2_photo")
    if not actor1_photo or not actor2_photo:
        raise HTTPException(status_code=400, detail="Wymagane dwa zdjęcia aktorów.")
    for photo in (actor1_photo, actor2_photo):
        if photo.content_type and photo.content_type not in allowed_img:
            raise HTTPException(status_code=415, detail="Wymagane zdjęcia JPG/PNG/WEBP.")

    voice_key_1 = form.get("voice_key_1", "pl_male_marek")
    voice_key_2 = form.get("voice_key_2", "pl_female_zofia")
    lines_raw = json.loads(form.get("lines", "[]"))
    if not lines_raw:
        raise HTTPException(status_code=400, detail="Skrypt jest pusty.")

    catalog = await get_full_catalog()

    def resolve_voice_id(key: str) -> str:
        info = catalog.get(key)
        if not info or info.get("type") == "cloned":
            return "pl-PL-MarekNeural"
        return info["id"]

    voice_id_1 = resolve_voice_id(voice_key_1)
    voice_id_2 = resolve_voice_id(voice_key_2)

    job_id = uuid.uuid4().hex[:8]
    work_dir = AUDIO_DIR / f"dialogue-{job_id}"
    work_dir.mkdir(parents=True, exist_ok=True)

    photo1_path = work_dir / f"actor1{Path(actor1_photo.filename or 'a.jpg').suffix}"
    photo2_path = work_dir / f"actor2{Path(actor2_photo.filename or 'a.jpg').suffix}"
    photo1_path.write_bytes(await actor1_photo.read())
    photo2_path.write_bytes(await actor2_photo.read())

    # Annotate each line with resolved voice_id
    lines = [
        {"actor": int(l.get("actor", 1)), "text": l["text"].strip(),
         "voice_id": voice_id_1 if int(l.get("actor", 1)) == 1 else voice_id_2}
        for l in lines_raw if l.get("text", "").strip()
    ]
    if not lines:
        raise HTTPException(status_code=400, detail="Brak kwestii z tekstem.")

    background_tasks.add_task(process_sadtalker_job, job_id, photo1_path, photo2_path, lines, work_dir)
    return {"job_id": job_id}


# ─── Presentation Studio ──────────────────────────────────────────────────────

_PRES_SESSIONS: Dict[str, dict] = {}


def _make_placeholder_slide(out_path: Path, slide_num: int, text: str):
    try:
        from PIL import Image, ImageDraw
        img = Image.new("RGB", (1280, 720), color=(24, 24, 40))
        draw = ImageDraw.Draw(img)
        draw.rectangle([30, 30, 1250, 690], outline=(80, 80, 140), width=2)
        draw.text((50, 50), f"Slajd {slide_num}", fill=(160, 160, 255))
        y = 110
        for line in text[:300].split("\n")[:12]:
            draw.text((50, y), line[:90], fill=(200, 200, 200))
            y += 40
        img.save(str(out_path))
    except Exception:
        # minimal 1x1 PNG fallback
        out_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00"
            b"\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )


def _extract_pdf_slides(pdf_path: Path, out_dir: Path) -> list:
    import fitz  # pymupdf
    doc = fitz.open(str(pdf_path))
    slides = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(2.5, 2.5)  # ~180 DPI — ostry tekst na slajdach po skalowaniu do 1080p
        pix = page.get_pixmap(matrix=mat)
        img_path = out_dir / f"slide_{i:03d}.png"
        pix.save(str(img_path))
        text = page.get_text().strip()
        slides.append({"index": i, "img_path": str(img_path), "slide_text": text, "notes": ""})
    doc.close()
    return slides


def _extract_pptx_notes(pptx_path: Path) -> list:
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    result = []
    for i, slide in enumerate(prs.slides):
        notes = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            notes = "\n".join(p.text for p in tf.paragraphs if p.text.strip())
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        result.append({"index": i, "notes": notes, "slide_text": "\n".join(texts)})
    return result


async def _pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Optional[Path]:
    proc = await asyncio.create_subprocess_exec(
        "libreoffice", "--headless", "--convert-to", "pdf",
        "--outdir", str(out_dir), str(pptx_path),
        stderr=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.DEVNULL,
    )
    await proc.communicate()
    if proc.returncode != 0:
        return None
    pdfs = list(out_dir.glob("*.pdf"))
    return pdfs[0] if pdfs else None


@app.post("/api/presentation/parse")
async def parse_presentation(file: UploadFile = File(...)):
    ext = Path(file.filename or "").suffix.lower()
    if ext not in (".pptx", ".pdf", ".ppt"):
        raise HTTPException(400, "Obsługiwane formaty: PPTX, PPT, PDF")

    session_id = uuid.uuid4().hex[:12]
    session_dir = AUDIO_DIR / f"pres-{session_id}"
    session_dir.mkdir(parents=True, exist_ok=True)

    file_path = session_dir / f"input{ext}"
    file_path.write_bytes(await file.read())

    try:
        if ext == ".pdf":
            slides = _extract_pdf_slides(file_path, session_dir)
        else:
            pptx_data = _extract_pptx_notes(file_path)
            pdf_path = await _pptx_to_pdf(file_path, session_dir)
            if pdf_path:
                img_slides = _extract_pdf_slides(pdf_path, session_dir)
                notes_map = {s["index"]: s for s in pptx_data}
                for s in img_slides:
                    d = notes_map.get(s["index"], {})
                    s["notes"] = d.get("notes", "")
                    s["slide_text"] = d.get("slide_text", s.get("slide_text", ""))
                slides = img_slides
            else:
                slides = []
                for d in pptx_data:
                    ph = session_dir / f"slide_{d['index']:03d}.png"
                    _make_placeholder_slide(ph, d["index"] + 1, d["slide_text"])
                    slides.append({**d, "img_path": str(ph)})

        import base64
        from PIL import Image as _PILImage
        import io as _io
        result = []
        for s in slides:
            img_path = Path(s["img_path"])
            img_data = ""
            if img_path.exists() and img_path.stat().st_size > 10:
                # Thumbnail JPEG (max 900px wide) for preview — full PNG kept for rendering
                try:
                    with _PILImage.open(img_path) as im:
                        w, h = im.size
                        if w > 900:
                            im = im.resize((900, int(h * 900 / w)), _PILImage.LANCZOS)
                        buf = _io.BytesIO()
                        im.convert("RGB").save(buf, format="JPEG", quality=80, optimize=True)
                        b64 = base64.b64encode(buf.getvalue()).decode()
                        img_data = f"data:image/jpeg;base64,{b64}"
                except Exception:
                    b64 = base64.b64encode(img_path.read_bytes()).decode()
                    img_data = f"data:image/png;base64,{b64}"
            result.append({
                "index": s["index"],
                "image": img_data,
                "notes": s.get("notes", ""),
                "slide_text": s.get("slide_text", ""),
            })

        _PRES_SESSIONS[session_id] = {
            "slides_dir": str(session_dir),
            "slides": {s["index"]: s for s in slides},
        }

        return {"session_id": session_id, "slides": result, "total": len(result), "filename": file.filename}

    except Exception as e:
        shutil.rmtree(str(session_dir), ignore_errors=True)
        raise HTTPException(500, f"Błąd parsowania: {e}")


class PresentationNarrationRequest(BaseModel):
    slide_text: str
    slide_num: int = 1
    total_slides: int = 1
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"


@app.post("/api/presentation/generate-narration")
async def generate_slide_narration(req: PresentationNarrationRequest):
    system = (
        "You are a professional presenter. Write clear, natural spoken narration for a slide. "
        "Use the same language as the slide text. Write 1-4 sentences that explain the slide content. "
        "Output ONLY the narration text — no JSON, no formatting, no commentary."
    )
    prompt = f"Slide {req.slide_num} of {req.total_slides}. Slide content:\n{req.slide_text}"
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.post(
            f"{req.ollama_host}/api/generate",
            json={"model": req.model, "prompt": prompt, "system": system, "stream": False, "keep_alive": 0},
        )
        data = res.json()
    return {"narration": data.get("response", "").strip()}


async def _process_presentation_job(job_id: str, session_id: str, edited_slides: list, voice_key: str, xtts_speed: float = 1.0):
    _total = len(edited_slides)
    JOBS[job_id] = {"status": "running", "progress": 0, "slides_done": 0, "slides_total": _total,
                    "message": "Inicjalizacja...", "eta_secs": None, "url": None, "error": None}
    session = _PRES_SESSIONS.get(session_id)
    if not session:
        JOBS[job_id] = {"status": "failed", "error": "Sesja wygasła — prześlij plik ponownie.", "url": None, "progress": None}
        return

    slides_dir = Path(session["slides_dir"])
    orig = session["slides"]
    tmp = slides_dir / f"render-{job_id}"
    tmp.mkdir(exist_ok=True)

    try:
        catalog = await get_full_catalog()
        voice_info = catalog.get(voice_key)

        async def make_tts(text: str, out: str):
            if voice_info and voice_info.get("type") == "cloned":
                ref = str(XTTS_SPEAKERS_DIR / voice_info["id"])
                voice_lang = voice_info.get("voice_lang") or _load_voice_lang(
                    voice_info["id"].rsplit(".", 1)[0]
                ) or "pl"
                try:
                    # Chatterbox: EN=Turbo, PL/other=Multilingual V3 (better quality, no CJK artifacts)
                    await chatterbox_synthesize(text, ref, out, language=voice_lang)
                except Exception as _cb_err:
                    logger.warning("Chatterbox unavailable (%s), falling back to XTTS", _cb_err)
                    await xtts_synthesize(text, ref, out, speed=xtts_speed, force_language=voice_lang)
            else:
                vid = (voice_info or {}).get("id", "pl-PL-MarekNeural")
                await edge_tts.Communicate(text, vid).save(out)

        clip_paths = []
        _slide_times: list = []
        for i, se in enumerate(edited_slides):
            notes = (se.get("notes") or "").strip()
            idx = se.get("index", i)
            _eta = int(sum(_slide_times) / len(_slide_times) * (_total - i)) if _slide_times else None
            JOBS[job_id].update({
                "progress": round(i / _total * 95) if _total else 0,
                "slides_done": i,
                "slides_total": _total,
                "message": f"Slajd {i + 1}/{_total}",
                "eta_secs": _eta,
            })
            _t0 = time.monotonic()
            img_path = (orig.get(idx) or {}).get("img_path", "")
            if not img_path or not Path(img_path).exists():
                _slide_times.append(time.monotonic() - _t0)
                continue

            vf = "scale=1920:1080:force_original_aspect_ratio=decrease:flags=lanczos,pad=1920:1080:(ow-iw)/2:(oh-ih)/2"
            out_clip = str(tmp / f"clip_{i:03d}.mp4")

            if not notes:
                # Include silent audio track so all clips are compatible for concat
                proc = await asyncio.create_subprocess_exec(
                    "ffmpeg", "-y",
                    "-loop", "1", "-i", img_path,
                    "-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo",
                    "-t", "3",
                    "-c:v", "libx264", "-tune", "stillimage", "-crf", "16", "-preset", "medium",
                    "-c:a", "aac", "-b:a", "192k",
                    "-pix_fmt", "yuv420p", "-vf", vf,
                    "-shortest", out_clip,
                    stderr=asyncio.subprocess.DEVNULL, stdout=asyncio.subprocess.DEVNULL,
                )
                await proc.communicate()
                if proc.returncode == 0:
                    clip_paths.append(out_clip)
                _slide_times.append(time.monotonic() - _t0)
                continue

            audio_path = str(tmp / f"audio_{i:03d}.mp3")
            await make_tts(notes, audio_path)

            # -shortest + apad don't work with -loop 1: -shortest terminates at the
            # unfiltered input duration, before apad adds its silence. Fix: set image
            # loop duration explicitly to audio_dur + pad so apad actually runs.
            _pad_s = 0.4
            _audio_dur_s = _ffmpeg_mp3_duration_ms(audio_path) / 1000.0
            _clip_dur_s = _audio_dur_s + _pad_s

            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-y",
                "-loop", "1", "-t", str(_clip_dur_s), "-i", img_path,
                "-i", audio_path,
                "-c:v", "libx264", "-tune", "stillimage", "-crf", "16", "-preset", "medium",
                "-c:a", "aac", "-b:a", "192k",
                "-pix_fmt", "yuv420p", "-vf", vf,
                "-af", "apad=pad_dur=0.4",
                out_clip,
                stderr=asyncio.subprocess.DEVNULL, stdout=asyncio.subprocess.DEVNULL,
            )
            await proc.communicate()
            if proc.returncode == 0:
                clip_paths.append(out_clip)
            _slide_times.append(time.monotonic() - _t0)

        if not clip_paths:
            raise RuntimeError("Nie wygenerowano żadnego klipu.")

        JOBS[job_id].update({"progress": 95, "message": "Scalanie klipów...", "eta_secs": None})
        list_file = tmp / "concat.txt"
        list_file.write_text("\n".join(f"file '{p}'" for p in clip_paths))
        _orig_name = Path(session.get("filename", f"presentation-{job_id}")).stem
        final_path = AUDIO_DIR / f"{_orig_name}-{job_id[:6]}.mp4"
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-y", "-f", "concat", "-safe", "0",
            "-i", str(list_file), "-c", "copy", str(final_path),
            stderr=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.DEVNULL,
        )
        _, err = await proc.communicate()
        if proc.returncode != 0:
            raise RuntimeError(f"ffmpeg concat: {err.decode()[-200:]}")

        JOBS[job_id] = {"status": "completed", "progress": 100,
                        "url": f"/api/audio/{final_path.name}",
                        "filename": final_path.name, "message": "Gotowe!", "mode": "video"}
    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e), "url": None, "progress": None}
    finally:
        shutil.rmtree(str(tmp), ignore_errors=True)


@app.post("/api/presentation/render")
async def render_presentation(request: Request, background_tasks: BackgroundTasks):
    _night_mode_block()
    form = await request.form()
    session_id = form.get("session_id", "")
    if session_id not in _PRES_SESSIONS:
        raise HTTPException(400, "Nieprawidłowe ID sesji — prześlij plik ponownie.")
    try:
        edited_slides = json.loads(form.get("slides", "[]"))
    except json.JSONDecodeError:
        raise HTTPException(400, "Nieprawidłowy format danych.")
    if not any(s.get("notes", "").strip() for s in edited_slides):
        raise HTTPException(400, "Żaden slajd nie ma tekstu narracji.")
    voice_key = form.get("voice_key", "pl_male_marek")
    xtts_speed = float(form.get("xtts_speed", "1.0"))
    job_id = uuid.uuid4().hex[:8]
    # Use job queue so multiple presentations are processed sequentially
    background_tasks.add_task(_enqueue_job, job_id, _process_presentation_job,
                              job_id, session_id, edited_slides, voice_key, xtts_speed)
    return {"job_id": job_id, "queued": True, "queue_length": len(_JOB_QUEUE) + 1}


def _secs_to_srt_time(s: float) -> str:
    h = int(s // 3600)
    m = int((s % 3600) // 60)
    sec = int(s % 60)
    ms = int((s % 1) * 1000)
    return f"{h:02d}:{m:02d}:{sec:02d},{ms:03d}"


@app.post("/api/edit-video")
async def edit_video(request: Request, background_tasks: BackgroundTasks):
    form = await request.form()
    video_upload: UploadFile = form.get("video")
    if not video_upload:
        raise HTTPException(400, "Brak pliku wideo.")

    trim_start_raw = form.get("trim_start")
    trim_end_raw = form.get("trim_end")
    trim_start = float(trim_start_raw) if trim_start_raw else 0.0
    trim_end = float(trim_end_raw) if trim_end_raw else -1.0
    audio_action = form.get("audio_action", "keep")
    tts_text = form.get("tts_text", "").strip()
    tts_speaker = form.get("tts_speaker", "pl_male_marek")
    subtitles_raw = form.get("subtitles", "")
    audio_upload: Optional[UploadFile] = form.get("audio_file")

    def _int_field(name: str) -> int:
        raw = form.get(name)
        try:
            return max(0, int(float(raw))) if raw else 0
        except (ValueError, TypeError):
            return 0

    crop_top = _int_field("crop_top")
    crop_bottom = _int_field("crop_bottom")
    crop_left = _int_field("crop_left")
    crop_right = _int_field("crop_right")
    rotate_deg = form.get("rotate_deg", "0")
    speed_raw = form.get("speed")
    speed = float(speed_raw) if speed_raw else 1.0

    tmp = Path(tempfile.mkdtemp(prefix="ve_"))
    try:
        ext = Path(video_upload.filename or "video.mp4").suffix.lower() or ".mp4"
        current = tmp / f"input{ext}"
        current.write_bytes(await video_upload.read())
        step = 0

        async def ffmpeg(*args):
            nonlocal step, current
            step += 1
            out = tmp / f"step{step}.mp4"
            cmd = ["ffmpeg", "-y", "-i", str(current), *[str(a) for a in args], str(out)]
            proc = await asyncio.create_subprocess_exec(
                *cmd, stderr=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.DEVNULL
            )
            _, err = await proc.communicate()
            if proc.returncode != 0:
                raise HTTPException(500, f"ffmpeg error: {err.decode()[-300:]}")
            current = out

        # 1. Trim
        if trim_start > 0 or trim_end > 0:
            args = []
            if trim_start > 0:
                args += ["-ss", str(trim_start)]
            if trim_end > 0:
                args += ["-to", str(trim_end)]
            await ffmpeg(*args, "-c", "copy")

        # 1b. Crop / rotate / speed — requires re-encoding the video stream
        vfilters = []
        if crop_top or crop_bottom or crop_left or crop_right:
            vfilters.append(
                f"crop=in_w-{crop_left + crop_right}:in_h-{crop_top + crop_bottom}:{crop_left}:{crop_top}"
            )
        if rotate_deg == "90":
            vfilters.append("transpose=1")
        elif rotate_deg == "180":
            vfilters.append("transpose=1,transpose=1")
        elif rotate_deg == "270":
            vfilters.append("transpose=2")
        afilters = []
        if speed and speed != 1.0:
            speed = max(0.5, min(2.0, speed))
            vfilters.append(f"setpts={1/speed}*PTS")
            afilters.append(f"atempo={speed}")

        if vfilters:
            args = ["-vf", ",".join(vfilters), "-c:v", "libx264", "-preset", "medium", "-crf", "18"]
            if afilters:
                args += ["-af", ",".join(afilters), "-c:a", "aac"]
            else:
                args += ["-c:a", "copy"]
            await ffmpeg(*args)

        # 2. Audio action
        if audio_action == "remove":
            await ffmpeg("-an", "-c:v", "copy")

        elif audio_action in ("replace_tts", "replace_file"):
            audio_path = tmp / "replacement.mp3"

            if audio_action == "replace_tts":
                if not tts_text:
                    raise HTTPException(400, "Brak tekstu TTS.")
                catalog = await get_full_catalog()
                voice_info = catalog.get(tts_speaker)
                if voice_info and voice_info.get("type") == "cloned":
                    ref_wav = str(XTTS_SPEAKERS_DIR / voice_info["id"])
                    await xtts_synthesize(tts_text, ref_wav, str(audio_path))
                else:
                    voice_id = (voice_info or {}).get("id", "pl-PL-MarekNeural")
                    communicate = edge_tts.Communicate(tts_text, voice_id)
                    await communicate.save(str(audio_path))
            else:
                if not audio_upload:
                    raise HTTPException(400, "Brak pliku audio.")
                audio_path.write_bytes(await audio_upload.read())

            step += 1
            out = tmp / f"step{step}.mp4"
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-y", "-i", str(current), "-i", str(audio_path),
                "-map", "0:v", "-map", "1:a", "-c:v", "copy", "-shortest", str(out),
                stderr=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.DEVNULL,
            )
            _, err = await proc.communicate()
            if proc.returncode != 0:
                raise HTTPException(500, f"ffmpeg replace audio error: {err.decode()[-300:]}")
            current = out

        # 3. Subtitles burn-in
        if subtitles_raw:
            try:
                subs = json.loads(subtitles_raw)
            except json.JSONDecodeError:
                subs = []
            if subs:
                srt_path = tmp / "subs.srt"
                lines_srt = []
                for i, sub in enumerate(subs, 1):
                    t_start = _secs_to_srt_time(float(sub.get("start", 0)))
                    t_end = _secs_to_srt_time(float(sub.get("end", float(sub.get("start", 0)) + 3)))
                    lines_srt.append(f"{i}\n{t_start} --> {t_end}\n{sub['text']}\n")
                srt_path.write_text("\n".join(lines_srt), encoding="utf-8")

                step += 1
                out = tmp / f"step{step}.mp4"
                proc = await asyncio.create_subprocess_exec(
                    "ffmpeg", "-y", "-i", str(current),
                    "-vf", f"subtitles={str(srt_path)}:force_style='FontSize=20,PrimaryColour=&H00FFFFFF,OutlineColour=&H00000000,BorderStyle=3'",
                    "-c:a", "copy", str(out),
                    stderr=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.DEVNULL,
                )
                _, err = await proc.communicate()
                if proc.returncode != 0:
                    raise HTTPException(500, f"ffmpeg subtitles error (sprawdź czy libass jest zainstalowany): {err.decode()[-300:]}")
                current = out

        result_bytes = current.read_bytes()
        out_name = f"edited_{Path(video_upload.filename or 'video.mp4').name}"

    finally:
        background_tasks.add_task(shutil.rmtree, str(tmp), True)

    from fastapi.responses import Response as FastResponse
    return FastResponse(
        content=result_bytes,
        media_type="video/mp4",
        headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
    )


# ─── Audiobook Studio ─────────────────────────────────────────────────────────

_AUDIOBOOK_SESSIONS: Dict[str, dict] = {}

# Max chars per single TTS call — prevents edge-tts truncation and XTTS OOM
_TTS_CHUNK_MAX = 200  # matches _XTTS_CHAR_LIMIT — safe for Polish XTTS hard limit of 224

# Max concurrent edge-tts requests
_TTS_SEMAPHORE = asyncio.Semaphore(6)


def _html_to_text(html: str) -> str:
    import re
    html = re.sub(r'<(script|style)[^>]*>.*?</(script|style)>', '', html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<br\s*/?>', '\n', html, flags=re.IGNORECASE)
    text = re.sub(r'</p>', '\n\n', text, flags=re.IGNORECASE)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _split_for_tts(text: str, max_chars: int = _TTS_CHUNK_MAX) -> list:
    """Split text into TTS-safe chunks at sentence/clause boundaries."""
    import re
    if len(text) <= max_chars:
        return [text]
    # Split at sentence endings first
    parts = re.split(r'(?<=[.!?…])\s+', text)
    chunks, current = [], ""
    for part in parts:
        if len(current) + len(part) + 1 <= max_chars:
            current = (current + " " + part).strip() if current else part
        else:
            if current:
                chunks.append(current)
            if len(part) > max_chars:
                # Force-split at comma/space
                for sub in re.split(r'(?<=,)\s+', part):
                    if len(sub) > max_chars:
                        for i in range(0, len(sub), max_chars):
                            chunks.append(sub[i:i + max_chars])
                    else:
                        chunks.append(sub)
                current = ""
            else:
                current = part
    if current:
        chunks.append(current)
    return [c for c in chunks if c.strip()]


def _extract_epub_chapters(path: Path) -> list:
    try:
        import ebooklib
        from ebooklib import epub
        import re
        book = epub.read_epub(str(path), options={'ignore_ncx': True})
        chapters = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            raw = item.get_content().decode('utf-8', errors='replace')
            text = _html_to_text(raw)
            if len(text) < 150:
                continue
            m = re.search(r'<h[1-3][^>]*>(.*?)</h[1-3]>', raw, re.IGNORECASE | re.DOTALL)
            title = re.sub(r'<[^>]+>', '', m.group(1)).strip() if m else item.get_name().split('/')[-1].rsplit('.', 1)[0]
            chapters.append({'title': title or f'Rozdział {len(chapters)+1}', 'text': text})
        return chapters
    except Exception as e:
        raise Exception(f"Błąd EPUB: {e}")


def _extract_txt_chapters(path: Path) -> list:
    import re
    text = path.read_text(encoding='utf-8', errors='replace')
    pattern = re.compile(
        r'\n{2,}((?:Rozdział|ROZDZIAŁ|Rozdz\.|Chapter|CHAPTER|Część|CZĘŚĆ)\s+[IVXLC\d]+[^\n]*)\n',
        re.IGNORECASE
    )
    parts = pattern.split(text)
    if len(parts) > 2:
        chapters = []
        if parts[0].strip():
            chapters.append({'title': 'Wstęp', 'text': parts[0].strip()})
        for i in range(1, len(parts), 2):
            body = parts[i + 1].strip() if i + 1 < len(parts) else ''
            if body:
                chapters.append({'title': parts[i].strip(), 'text': body})
        return chapters
    paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip()]
    chunks, current, cur_len, num = [], [], 0, 1
    for p in paragraphs:
        if cur_len + len(p) > 5000 and current:
            chunks.append({'title': f'Część {num}', 'text': '\n\n'.join(current)})
            num += 1; current = [p]; cur_len = len(p)
        else:
            current.append(p); cur_len += len(p)
    if current:
        chunks.append({'title': f'Część {num}', 'text': '\n\n'.join(current)})
    return chunks


def _extract_pdf_chapters_audiobook(path: Path) -> list:
    import fitz
    doc = fitz.open(str(path))
    pages = [page.get_text().strip() for page in doc if page.get_text().strip()]
    doc.close()
    chunks, current, cur_len, num = [], [], 0, 1
    for pt in pages:
        if cur_len + len(pt) > 4500 and current:
            chunks.append({'title': f'Część {num}', 'text': '\n\n'.join(current)})
            num += 1; current = [pt]; cur_len = len(pt)
        else:
            current.append(pt); cur_len += len(pt)
    if current:
        chunks.append({'title': f'Część {num}', 'text': '\n\n'.join(current)})
    return chunks


def _load_xtts_gpu_sync(models_dir: str):
    import os as _os
    _os.environ.setdefault("TTS_AGREE_TO_USER_AGREEMENT", "1")
    _os.environ.setdefault("COQUI_TOS_AGREED", "1")
    from TTS.api import TTS as CoquiTTS
    return CoquiTTS("tts_models/multilingual/multi-dataset/xtts_v2", gpu=True)


async def _detect_characters_llm(sample: str, ollama_host: str, model: str) -> list:
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            res = await client.post(f"{ollama_host}/api/generate", json={
                "model": model, "stream": False,
                "prompt": (
                    "List the main speaking characters in this book excerpt (not the narrator). "
                    "Return ONLY JSON: {\"characters\": [\"Name1\", \"Name2\"]} "
                    "Max 10 names. Only characters who have dialogue lines.\n\n"
                    f"Text:\n{sample[:2500]}"
                ),
                "keep_alive": "5m",
            })
            result = json.loads(res.json().get("response", "{}"))
            return [c for c in result.get("characters", []) if isinstance(c, str) and c.strip()]
    except Exception:
        return []


def _segment_chapter_regex(text: str, characters: list) -> list:
    """Fast regex-based segmentation — no LLM, handles big chapters instantly."""
    import re

    # Build character name matcher (case-insensitive)
    char_set = {c.lower() for c in characters}
    attrib_re = re.compile(
        r'(?:—\s*)?(?:powiedział|rzekł|odparł|mruknął|krzyknął|wyszeptał|zapytał|odrzekł|stwierdzić|stwierdził|dodał|przerwał|szepnął|zawołał|warknął|syknął)\s+(\w+)',
        re.IGNORECASE
    )

    segments = []
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    last_speaker: Optional[str] = None

    for line in lines:
        is_dialogue = line.startswith('—') or line.startswith('—')
        is_epigraph = len(line) < 120 and (line.startswith('"') or line.startswith('“'))

        if is_epigraph and not segments:
            seg_type = "epigraph"
            character = None
            emotion = "formal"
        elif is_dialogue:
            seg_type = "dialogue"
            m = attrib_re.search(line)
            character = None
            if m:
                name = m.group(1)
                if name.lower() in char_set:
                    character = name.capitalize()
                    last_speaker = character
            if character is None and last_speaker:
                character = last_speaker
            emotion = "dramatic" if '!' in line else ("suspense" if ('...' in line or '…' in line) else "calm")
        else:
            seg_type = "narrator"
            character = None
            last_speaker = None  # reset speaker after narration block
            lw = line.lower()
            if any(w in lw for w in ['biegł', 'walka', 'atak', 'szybko', 'nagle', 'błyskawicznie', 'rzucił się', 'skoczył']):
                emotion = "action"
            elif any(w in lw for w in ['cisza', 'mrok', 'ciemność', 'strach', 'groza', 'niebezpiecz', 'drzał']):
                emotion = "suspense"
            else:
                emotion = "calm"

        for chunk in _split_for_tts(line):
            segments.append({"type": seg_type, "character": character, "emotion": emotion, "text": chunk})

    return segments


async def _segment_chapter_llm(text: str, characters: list, ollama_host: str, model: str) -> list:
    """LLM-based segmentation for shorter chapters (<8K chars). Falls back to regex."""
    if len(text) > 8000:
        # Too long for LLM — use regex directly
        return _segment_chapter_regex(text, characters)

    char_list = ', '.join(characters[:10]) if characters else 'none detected'
    system = (
        f"Segment this book excerpt. Known characters: {char_list}.\n"
        "For each segment: type (narrator/dialogue/internal/epigraph), "
        "character (name or null), emotion (calm/action/suspense/dramatic/formal), text.\n"
        "Keep segments SHORT (1-2 sentences max). Polish dialogue starts with —.\n"
        'Return ONLY JSON: {"segments": [...]}'
    )
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(f"{ollama_host}/api/chat", json={
                "model": model, "stream": False,
                "messages": [{"role": "system", "content": system}, {"role": "user", "content": text}],
                "keep_alive": "5m",
            })
            result = _extract_json(res.json().get("message", {}).get("content", "{}"))
            segs = result.get("segments", [])
            if segs:
                # Ensure each segment respects TTS limit
                final = []
                for s in segs:
                    for chunk in _split_for_tts(s.get("text", "")):
                        final.append({**s, "text": chunk})
                return final
    except Exception:
        pass
    return _segment_chapter_regex(text, characters)


@app.post("/api/audiobook/parse")
async def parse_audiobook(
    file: UploadFile = File(...),
    ollama_host: str = Form("http://localhost:11434"),
    model: str = Form("qwen3.5:4b"),
):
    ext = Path(file.filename or "").suffix.lower()
    if ext not in (".epub", ".pdf", ".txt"):
        raise HTTPException(400, "Obsługiwane formaty: EPUB, PDF, TXT")

    session_id = uuid.uuid4().hex[:12]
    session_dir = AUDIO_DIR / f"ab-parse-{session_id}"
    session_dir.mkdir(parents=True, exist_ok=True)
    file_path = session_dir / f"input{ext}"
    file_path.write_bytes(await file.read())

    try:
        if ext == ".epub":
            chapters = _extract_epub_chapters(file_path)
        elif ext == ".pdf":
            chapters = _extract_pdf_chapters_audiobook(file_path)
        else:
            chapters = _extract_txt_chapters(file_path)

        if not chapters:
            raise Exception("Nie udało się wyodrębnić tekstu z pliku.")

        sample = ' '.join(c['text'] for c in chapters[:2])[:3000]
        characters = await _detect_characters_llm(sample, ollama_host, model)

        _AUDIOBOOK_SESSIONS[session_id] = {
            "chapters": chapters,
            "characters": characters,
        }

        return {
            "session_id": session_id,
            "title": Path(file.filename or "audiobook").stem,
            "total_chapters": len(chapters),
            "characters": characters,
            "chapters_preview": [
                {"index": i, "title": c["title"], "preview": c["text"][:200]}
                for i, c in enumerate(chapters[:6])
            ],
        }
    except Exception as e:
        shutil.rmtree(str(session_dir), ignore_errors=True)
        raise HTTPException(500, f"Błąd parsowania: {e}")


class AudiobookGenerateRequest(BaseModel):
    session_id: str
    voice_map: Dict[str, str]
    book_title: str = "audiobook"
    output_format: str = "mp3"
    use_gpu_xtts: bool = True
    use_llm_analysis: bool = True
    ollama_host: str = "http://localhost:11434"
    model: str = "qwen3.5:4b"
    xtts_speed: float = 1.0


_EMOTION_PROSODY: Dict[str, Dict[str, str]] = {
    "calm":     {"rate": "-5%",  "pitch": "+0Hz"},
    "action":   {"rate": "+18%", "pitch": "+5Hz"},
    "suspense": {"rate": "-18%", "pitch": "-10Hz"},
    "dramatic": {"rate": "-8%",  "pitch": "-8Hz"},
    "formal":   {"rate": "-12%", "pitch": "-5Hz"},
}


async def _tts_segment(
    seg_idx: int,
    text: str,
    voice_id: str,
    prosody: dict,
    seg_path: str,
    xtts_model,
    speaker_wav: Optional[str],
    xtts_speed: float = 1.0,
) -> bool:
    """Generate one TTS segment. Returns True on success. Respects global semaphore for edge-tts."""
    try:
        if xtts_model and speaker_wav:
            wav_path = seg_path.replace(".mp3", ".wav")
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(
                _xtts_executor, _xtts_synthesize_sync, xtts_model, text, speaker_wav, wav_path,
                _detect_xtts_lang(text), xtts_speed
            )
            _ffmpeg_wav_to_mp3(wav_path, seg_path)
            if os.path.exists(wav_path):
                os.unlink(wav_path)
        else:
            async with _TTS_SEMAPHORE:
                await edge_tts.Communicate(
                    text=text, voice=voice_id,
                    rate=prosody["rate"], pitch=prosody["pitch"],
                ).save(seg_path)
        return os.path.exists(seg_path) and os.path.getsize(seg_path) > 100
    except Exception as e:
        logger.warning("[audiobook] seg %d TTS error: %s", seg_idx, e)
        return False


async def process_audiobook_job(
    job_id: str,
    session_id: str,
    voice_map: dict,
    book_title: str,
    output_format: str,
    use_gpu_xtts: bool,
    use_llm_analysis: bool,
    ollama_host: str,
    model: str,
    xtts_speed: float = 1.0,
):
    JOBS[job_id] = {
        "status": "processing", "progress": 0, "mode": "audiobook",
        "message": "Inicjalizacja...", "chapters_done": 0, "chapters_total": 0,
        "chapter_urls": [],
    }
    session = _AUDIOBOOK_SESSIONS.get(session_id)
    if not session:
        JOBS[job_id] = {"status": "failed", "error": "Sesja wygasła — prześlij plik ponownie."}
        return

    chapters = session["chapters"]
    characters = session["characters"]
    narrator_voice = voice_map.get("narrator", "pl_male_andrew_multi")
    safe_title = _sanitize_filename(book_title)
    book_dir = AUDIO_DIR / safe_title
    book_dir.mkdir(parents=True, exist_ok=True)
    out_dir = book_dir / f".tmp-{job_id}"
    out_dir.mkdir(parents=True, exist_ok=True)

    try:
        JOBS[job_id]["chapters_total"] = len(chapters)

        # ── Phase 1: Segmentation ──────────────────────────────────────────────
        analyzed = []
        if use_llm_analysis:
            JOBS[job_id]["message"] = "Faza 1/3: Segmentacja tekstu (AI dla krótkich rozdziałów, regex dla długich)..."
            for i, ch in enumerate(chapters):
                JOBS[job_id]["progress"] = int(i / len(chapters) * 20)
                JOBS[job_id]["message"] = f"Segmentacja {i+1}/{len(chapters)}: {ch['title'][:28]}..."
                # LLM only for short chapters; regex for long ones (fast and reliable)
                segs = await _segment_chapter_llm(ch["text"], characters, ollama_host, model)
                analyzed.append({"title": ch["title"], "segments": segs})
            # Unload Ollama
            try:
                async with httpx.AsyncClient(timeout=8.0) as client:
                    await client.post(f"{ollama_host}/api/generate",
                                      json={"model": model, "keep_alive": 0})
            except Exception:
                pass
        else:
            JOBS[job_id]["message"] = "Faza 1/3: Segmentacja (regex)..."
            for ch in chapters:
                analyzed.append({"title": ch["title"],
                                  "segments": _segment_chapter_regex(ch["text"], characters)})

        # ── Phase 2: Load XTTS on GPU (only if cloned voices are in use) ───────
        cloned_voices_used = any(
            v for v in voice_map.values()
            if v.startswith("cloned_")
        )
        xtts_gpu = None
        if use_gpu_xtts and _check_gpu() and cloned_voices_used:
            JOBS[job_id]["message"] = "Faza 2/3: Ładowanie XTTS-v2 na GPU..."
            try:
                import torch
                torch.cuda.empty_cache()
                xtts_gpu = await asyncio.get_event_loop().run_in_executor(
                    _xtts_executor, _load_xtts_gpu_sync, str(MODELS_DIR)
                )
            except Exception as e:
                logger.warning("[audiobook] GPU XTTS error: %s — edge-tts only", e)
                xtts_gpu = None

        # ── Phase 3: Generate audio per chapter ────────────────────────────────
        JOBS[job_id]["message"] = "Faza 3/3: Generowanie audio..."
        total_segs = max(sum(len(c["segments"]) for c in analyzed), 1)
        segs_done = 0
        chapter_files: list = []  # [(title, path)]
        full_catalog = await get_full_catalog()

        for ch_idx, chapter in enumerate(analyzed):
            JOBS[job_id]["chapters_done"] = ch_idx
            JOBS[job_id]["message"] = (
                f"Rozdział {ch_idx+1}/{len(analyzed)}: {chapter['title'][:28]}... "
                f"({len(chapter['segments'])} segmentów)"
            )
            tmp = out_dir / f"ch_{ch_idx:03d}_tmp"
            tmp.mkdir(exist_ok=True)

            # Build task list for this chapter
            tasks_meta = []  # (seg_idx, seg_path, adds_pause)
            for seg_idx, seg in enumerate(chapter["segments"]):
                text = (seg.get("text") or "").strip()
                if len(text) < 2:
                    continue

                seg_type = seg.get("type", "narrator")
                character = seg.get("character")
                emotion = seg.get("emotion", "calm")
                prosody = _EMOTION_PROSODY.get(emotion, _EMOTION_PROSODY["calm"])

                # Resolve voice
                if seg_type == "dialogue" and character and character in voice_map:
                    vkey = voice_map[character]
                elif seg_type == "internal":
                    vkey = voice_map.get("internal", narrator_voice)
                elif seg_type == "epigraph":
                    vkey = voice_map.get("epigraph", narrator_voice)
                else:
                    vkey = narrator_voice

                voice_info = full_catalog.get(vkey)
                is_cloned = voice_info and voice_info.get("type") == "cloned"
                speaker_wav = str(XTTS_SPEAKERS_DIR / voice_info["id"]) if is_cloned else None
                voice_id = (voice_info or {}).get("id", "pl-PL-MarekNeural")

                seg_path = str(tmp / f"seg_{seg_idx:04d}.mp3")
                tasks_meta.append((seg_idx, text, voice_id, prosody, seg_path,
                                   xtts_gpu if is_cloned else None, speaker_wav,
                                   seg_type in ("dialogue", "epigraph")))

            # Run edge-tts tasks concurrently; XTTS tasks are sequential (GPU)
            edge_tasks = [(m[0], m[1], m[2], m[3], m[4], None, None, m[7])
                          for m in tasks_meta if m[5] is None]
            xtts_tasks = [m for m in tasks_meta if m[5] is not None]

            results: dict = {}  # seg_idx -> (path, adds_pause)

            # XTTS — sequential (single GPU)
            for seg_idx, text, voice_id, prosody, seg_path, model_obj, speaker_wav, adds_pause in xtts_tasks:
                ok = await _tts_segment(seg_idx, text, voice_id, prosody, seg_path, model_obj, speaker_wav, xtts_speed)
                if ok:
                    results[seg_idx] = (seg_path, adds_pause)
                segs_done += 1
                JOBS[job_id]["progress"] = 20 + int(segs_done / total_segs * 75)

            # edge-tts — parallel batches
            async def _run_edge(item):
                seg_idx, text, voice_id, prosody, seg_path, _, __, adds_pause = item
                ok = await _tts_segment(seg_idx, text, voice_id, prosody, seg_path, None, None)
                return seg_idx, seg_path, adds_pause, ok

            BATCH = 12  # concurrent edge-tts calls
            for batch_start in range(0, len(edge_tasks), BATCH):
                batch = edge_tasks[batch_start:batch_start + BATCH]
                batch_results = await asyncio.gather(*[_run_edge(item) for item in batch])
                for seg_idx, seg_path, adds_pause, ok in batch_results:
                    if ok:
                        results[seg_idx] = (seg_path, adds_pause)
                segs_done += len(batch)
                JOBS[job_id]["progress"] = 20 + int(segs_done / total_segs * 75)

            # Rebuild ordered file list
            seg_files = []
            silence_tmp = str(tmp)
            for seg_idx in sorted(results.keys()):
                seg_path, adds_pause = results[seg_idx]
                seg_files.append(seg_path)
                if adds_pause:
                    seg_files.append(await create_silence(150, silence_tmp))

            if seg_files:
                safe_ch = _sanitize_filename(chapter["title"])
                ch_name = f"{ch_idx+1:02d} - {safe_ch}.mp3"
                ch_raw = str(book_dir / f".raw_{ch_idx:03d}.mp3")
                ch_out = str(book_dir / ch_name)
                await concatenate_audio(seg_files, ch_raw)
                JOBS[job_id]["message"] = f"Mastering rozdziału {ch_idx+1}…"
                await _audiobook_master_chapter(ch_raw, ch_out)
                try:
                    os.unlink(ch_raw)
                except OSError:
                    pass
                chapter_files.append((chapter["title"], ch_out))
                # ── Incremental update — chapters downloadable as they finish ──
                JOBS[job_id]["chapter_urls"] = [
                    {"title": t, "url": f"/api/audio/{safe_title}/{os.path.basename(f)}"}
                    for t, f in chapter_files
                ]

            shutil.rmtree(str(tmp), ignore_errors=True)

        if xtts_gpu is not None:
            try:
                import torch
                del xtts_gpu
                torch.cuda.empty_cache()
            except Exception:
                pass

        if not chapter_files:
            raise Exception("Nie wygenerowano żadnego audio.")

        JOBS[job_id].update({"progress": 97, "message": "Scalanie rozdziałów..."})

        final_name = f"{safe_title}.mp3"
        final_path = book_dir / final_name
        list_path = out_dir / "concat.txt"
        list_path.write_text("\n".join(f"file '{f}'" for _, f in chapter_files))
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-f", "concat", "-safe", "0", "-i", str(list_path),
            "-c", "copy", str(final_path), "-y",
            stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
        )
        await proc.communicate()

        m4b_url = None
        if output_format == "m4b" and final_path.exists():
            m4b_name = f"{safe_title}.m4b"
            m4b_path = book_dir / m4b_name
            meta_lines = [";FFMETADATA1\n"]
            current_ms = 0
            for title, af in chapter_files:
                try:
                    dur_ms = _ffmpeg_mp3_duration_ms(af)
                except Exception:
                    dur_ms = 60000
                meta_lines.append(
                    f"\n[CHAPTER]\nTIMEBASE=1/1000\nSTART={current_ms}\n"
                    f"END={current_ms + dur_ms}\ntitle={title}\n"
                )
                current_ms += dur_ms
            meta_path = out_dir / "chapters.ffmeta"
            meta_path.write_text("".join(meta_lines))
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-i", str(final_path), "-i", str(meta_path),
                "-map_metadata", "1", "-c:a", "aac", "-b:a", "128k",
                "-movflags", "+faststart", str(m4b_path), "-y",
                stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
            )
            await proc.communicate()
            if m4b_path.exists():
                m4b_url = f"/api/audio/{safe_title}/{m4b_name}"

        shutil.rmtree(str(out_dir), ignore_errors=True)

        JOBS[job_id] = {
            "status": "completed", "progress": 100, "mode": "audiobook",
            "url": f"/api/audio/{safe_title}/{final_name}",
            "m4b_url": m4b_url,
            "filename": final_name,
            "folder": safe_title,
            "chapter_urls": [
                {"title": t, "url": f"/api/audio/{safe_title}/{os.path.basename(f)}"}
                for t, f in chapter_files
            ],
            "message": "Gotowe!",
        }

    except Exception as e:
        JOBS[job_id] = {"status": "failed", "error": str(e)}
        shutil.rmtree(str(out_dir), ignore_errors=True)


@app.post("/api/audiobook/generate")
async def generate_audiobook(req: AudiobookGenerateRequest, background_tasks: BackgroundTasks):
    _night_mode_block()
    if req.session_id not in _AUDIOBOOK_SESSIONS:
        raise HTTPException(400, "Nieprawidłowe session_id — prześlij plik ponownie.")
    job_id = uuid.uuid4().hex[:8]
    background_tasks.add_task(
        _enqueue_job, job_id, process_audiobook_job,
        job_id, req.session_id, req.voice_map, req.book_title,
        req.output_format, req.use_gpu_xtts, req.use_llm_analysis,
        req.ollama_host, req.model, req.xtts_speed,
    )
    return {"job_id": job_id, "queued": True, "queue_length": len(_JOB_QUEUE) + 1}


# ── Queue & Night Mode endpoints ─────────────────────────────────────────────

@app.get("/api/queue/status")
async def get_queue_status():
    queued_jobs = [
        {"job_id": jid, "position": JOBS.get(jid, {}).get("queue_position"), "status": "queued",
         "created_ts": JOBS.get(jid, {}).get("created_ts")}
        for jid, *_ in _JOB_QUEUE
    ]
    active = None
    if _JOB_QUEUE_ACTIVE and _JOB_QUEUE_ACTIVE in JOBS:
        active = {"job_id": _JOB_QUEUE_ACTIVE, **JOBS[_JOB_QUEUE_ACTIVE]}
    return {
        "running": _JOB_QUEUE_RUNNING,
        "queue_length": len(_JOB_QUEUE),
        "active_job": _JOB_QUEUE_ACTIVE,
        "active_job_info": active,
        "queued_jobs": queued_jobs,
    }

@app.delete("/api/queue/{job_id}")
async def cancel_queued_job(job_id: str):
    global _JOB_QUEUE
    if job_id == _JOB_QUEUE_ACTIVE:
        raise HTTPException(400, "Nie można anulować aktywnego zadania")
    before = len(_JOB_QUEUE)
    _JOB_QUEUE = deque(j for j in _JOB_QUEUE if j[0] != job_id)
    if job_id in JOBS:
        JOBS[job_id] = {"status": "cancelled", "progress": None, "url": None, "error": "Anulowane przez użytkownika"}
    # Renumber positions
    for i, (jid, *_) in enumerate(_JOB_QUEUE):
        if jid in JOBS:
            JOBS[jid]["queue_position"] = i + 1
    return {"cancelled": len(_JOB_QUEUE) < before, "queue_length": len(_JOB_QUEUE)}

class _NightModeReq(BaseModel):
    enabled: bool
    start_hour: int = 22
    end_hour: int = 7

@app.get("/api/settings/night-mode")
async def get_night_mode():
    return {**_NIGHT_MODE, "active_now": _is_night_mode_active()}

@app.post("/api/settings/night-mode")
async def set_night_mode(req: _NightModeReq):
    _NIGHT_MODE["enabled"] = req.enabled
    _NIGHT_MODE["start_hour"] = req.start_hour
    _NIGHT_MODE["end_hour"] = req.end_hour
    return {**_NIGHT_MODE, "active_now": _is_night_mode_active()}


# ── QA — Quality Assurance ────────────────────────────────────────────────────

async def _qa_run_cmd(*args: str) -> tuple[str, str]:
    proc = await asyncio.create_subprocess_exec(
        *args, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    stdout, stderr = await proc.communicate()
    return stdout.decode(errors="replace"), stderr.decode(errors="replace")

async def _qa_check_media(path: str) -> dict:
    p = Path(path)
    if not p.exists():
        return {"ok": False, "issues": ["file_not_found"], "warnings": [], "metrics": {}}
    size_kb = p.stat().st_size / 1024
    min_size = 10 if path.endswith(".mp4") else 1
    if size_kb < min_size:
        return {"ok": False, "issues": ["file_too_small"], "warnings": [], "metrics": {"size_kb": size_kb}}

    metrics: dict = {"size_kb": round(size_kb, 1)}
    issues: list = []
    warnings: list = []

    # ffprobe for streams and duration
    out, _ = await _qa_run_cmd(
        "ffprobe", "-v", "error",
        "-show_entries", "format=duration:stream=codec_type,sample_rate",
        "-of", "json", path,
    )
    try:
        probe = json.loads(out)
        streams = probe.get("streams", [])
        metrics["duration"] = float(probe.get("format", {}).get("duration", 0))
        metrics["has_video"] = any(s.get("codec_type") == "video" for s in streams)
        metrics["has_audio"] = any(s.get("codec_type") == "audio" for s in streams)
        metrics["sample_rate"] = next(
            (int(s["sample_rate"]) for s in streams if "sample_rate" in s), None
        )
    except (json.JSONDecodeError, ValueError):
        issues.append("ffprobe_failed")
        return {"ok": False, "issues": issues, "warnings": warnings, "metrics": metrics}

    is_video = path.lower().endswith(".mp4") or path.lower().endswith(".webm")
    min_dur = 1.0 if is_video else 0.5
    if metrics["duration"] < min_dur:
        issues.append(f"duration_too_short ({metrics['duration']:.2f}s)")
    if is_video and not metrics["has_audio"]:
        issues.append("missing_audio_stream")
    if is_video and not metrics["has_video"]:
        issues.append("missing_video_stream")

    # Volume check (audio files and video with audio)
    if metrics["has_audio"]:
        _, err = await _qa_run_cmd(
            "ffmpeg", "-i", path, "-af", "volumedetect", "-f", "null", "/dev/null",
        )
        m_vol = re.search(r"mean_volume:\s*([-\d.]+)", err)
        x_vol = re.search(r"max_volume:\s*([-\d.]+)", err)
        if m_vol:
            metrics["mean_volume_db"] = float(m_vol.group(1))
            if metrics["mean_volume_db"] < -60:
                issues.append("silence_detected")
        if x_vol:
            metrics["max_volume_db"] = float(x_vol.group(1))
            if metrics["max_volume_db"] >= -0.1:
                warnings.append("clipping_risk")

    return {"ok": len(issues) == 0, "issues": issues, "warnings": warnings, "metrics": metrics}

async def _qa_auto_fix_audio(path: str, issues: list) -> Optional[str]:
    if "silence_detected" in issues:
        return None  # cannot fix — need regeneration
    p = Path(path)
    fixed = p.parent / f"{p.stem}_fixed{p.suffix}"
    proc = await asyncio.create_subprocess_exec(
        "ffmpeg", "-y", "-i", path, "-af", "loudnorm=I=-16:TP=-1.5:LRA=11", str(fixed),
        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL,
    )
    await proc.communicate()
    if fixed.exists() and fixed.stat().st_size > 1024:
        return str(fixed)
    return None

class _QARequest(BaseModel):
    url: str  # e.g. /api/audio/filename.mp3 or absolute path

def _url_to_path(url: str) -> str:
    if url.startswith("/api/audio/"):
        fname = url.removeprefix("/api/audio/")
        # Use _safe_audio_path to support subdirectory paths (e.g. audiobook/chapter.mp3)
        # while still preventing path traversal attacks.
        return str(_safe_audio_path(fname))
    if Path(url).exists():
        return url
    raise HTTPException(400, "Nieprawidłowy URL. Użyj /api/audio/<nazwa_pliku>")

@app.post("/api/qa/check")
async def qa_check(req: _QARequest):
    path = _url_to_path(req.url)
    result = await _qa_check_media(path)
    result["file_url"] = req.url
    result["filename"] = Path(path).name
    return result

@app.post("/api/qa/fix")
async def qa_fix(req: _QARequest):
    path = _url_to_path(req.url)
    qa = await _qa_check_media(path)
    if qa["ok"] and not qa["warnings"]:
        return {"fixed": False, "message": "Plik nie wymaga naprawy", "url": req.url}
    fixed_path = await _qa_auto_fix_audio(path, qa["issues"] + qa["warnings"])
    if fixed_path is None:
        return {"fixed": False, "message": "Nie można auto-naprawić — plik zawiera tylko ciszę, wymagana regeneracja", "url": req.url}
    fixed_name = Path(fixed_path).name
    return {
        "fixed": True,
        "message": "loudnorm zastosowany",
        "url": f"/api/audio/{fixed_name}",
        "filename": fixed_name,
        "original_url": req.url,
    }

# ── FAZA 7: Advanced TTS stubs ────────────────────────────────────────────────

class _CosyVoiceReq(BaseModel):
    text: str
    reference_audio_url: str
    language: str = "zh"

class _OpenVoiceReq(BaseModel):
    text: str
    reference_audio_url: str
    style: str = "default"

@app.post("/api/tts/cosyvoice")
async def cosyvoice_tts(req: _CosyVoiceReq):
    try:
        import cosyvoice as _cv  # noqa: F401
    except ImportError:
        raise HTTPException(503, detail={"error": "cosyvoice_not_installed", "install_cmd": "pip install cosyvoice"})
    raise HTTPException(501, "CosyVoice integration pending GPU testing")

@app.post("/api/tts/openvoice")
async def openvoice_tts(req: _OpenVoiceReq):
    try:
        from openvoice import se_extractor as _se  # noqa: F401
    except ImportError:
        raise HTTPException(503, detail={"error": "openvoice_not_installed", "install_cmd": "git clone https://github.com/myshell-ai/OpenVoice && pip install -e OpenVoice/"})
    raise HTTPException(501, "OpenVoice integration pending")

_RVC_VENV = Path("/opt/rvc-venv")
_RVC_PYTHON = _RVC_VENV / "bin" / "python3"
_RVC_MODELS_DIR = Path.home() / "rvc-models"


def _check_rvc() -> bool:
    return _RVC_PYTHON.exists()


@app.post("/api/audio/rvc")
async def rvc_convert(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    model: str = Form(""),
    semitone_shift: int = Form(0),
):
    if not _check_rvc():
        raise HTTPException(503, detail={
            "error": "rvc_not_installed",
            "message": "RVC wymaga odizolowanego środowiska. Uruchom: python install.py",
            "install_cmd": "python3 -m venv /opt/rvc-venv && /opt/rvc-venv/bin/pip install rvc-python 'numpy<2' torchcrepe fairseq ffmpeg-python",
        })

    if model and not re.match(r"^[a-zA-Z0-9_\-\.]{1,128}$", model):
        raise HTTPException(400, detail="Nieprawidłowa nazwa modelu RVC.")
    model_path = _RVC_MODELS_DIR / model if model else None
    if model_path and not model_path.exists():
        raise HTTPException(404, detail=f"Model RVC '{model}' nie znaleziony w {_RVC_MODELS_DIR}")

    job_id = uuid.uuid4().hex[:8]
    JOBS[job_id] = {"status": "running", "progress": 10, "message": "Konwersja głosu (RVC)…", "url": None}

    suffix = Path(file.filename or "audio.wav").suffix or ".wav"
    tmp_dir = Path(tempfile.mkdtemp(prefix="rvc_"))
    in_path = tmp_dir / f"input{suffix}"
    in_path.write_bytes(await file.read())

    async def _run():
        try:
            out_path = tmp_dir / f"rvc_out_{job_id}.wav"
            # Pass all paths via JSON config to prevent injection from path names with quotes
            rvc_config = {
                "in_path": str(in_path),
                "out_path": str(out_path),
                "model_path": str(model_path) if model_path else None,
                "semitone_shift": semitone_shift,
                "site_packages": f"{_RVC_VENV}/lib/python3.11/site-packages",
            }
            config_path = tmp_dir / "rvc_config.json"
            config_path.write_text(json.dumps(rvc_config))
            script = """import sys, json
with open(sys.argv[1]) as _f:
    _c = json.load(_f)
sys.path.insert(0, _c["site_packages"])
from rvc_python.infer import RVCInference
_rvc = RVCInference(device="cuda:0")
if _c["model_path"]:
    _rvc.load_model(_c["model_path"])
_rvc.infer_file(_c["in_path"], _c["out_path"], f0_up_key=_c["semitone_shift"])
"""
            script_path = tmp_dir / "rvc_run.py"
            script_path.write_text(script)
            proc = await asyncio.create_subprocess_exec(
                str(_RVC_PYTHON), str(script_path), str(config_path),
                stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await proc.communicate()
            if proc.returncode != 0 or not out_path.exists():
                raise Exception(f"RVC błąd: {stderr.decode(errors='replace')[-400:]}")
            dest = AUDIO_DIR / f"rvc-{job_id}.wav"
            shutil.copy(str(out_path), str(dest))
            JOBS[job_id] = {"status": "completed", "progress": 100, "url": f"/api/audio/{dest.name}", "filename": dest.name}
        except Exception as e:
            JOBS[job_id] = {"status": "failed", "error": str(e), "url": None}
        finally:
            shutil.rmtree(str(tmp_dir), ignore_errors=True)

    background_tasks.add_task(_run)
    return {"job_id": job_id}

@app.post("/api/audio/demucs")
async def demucs_separate(background_tasks: BackgroundTasks, file: UploadFile = File(...), stems: str = "4"):
    """Separate audio into stems using Demucs.
    stems: '4' → vocals/drums/bass/other, '2' → vocals/accompaniment
    """
    try:
        import demucs  # noqa: F401
    except ImportError:
        raise HTTPException(503, detail={"error": "demucs_not_installed", "install_cmd": "pip install demucs"})
    job_id = uuid.uuid4().hex[:8]
    JOBS[job_id] = {"status": "running", "progress": 5, "message": "Separowanie ścieżek…", "url": None, "error": None}
    tmp = Path(tempfile.mkdtemp(prefix="demucs_"))
    in_path = tmp / f"input{Path(file.filename or 'audio.mp3').suffix}"
    in_path.write_bytes(await file.read())

    async def _run():
        try:
            JOBS[job_id]["message"] = "Uruchamianie Demucs (htdemucs)…"
            cmd = ["python3", "-m", "demucs", "--mp3", "-o", str(tmp)]
            if stems == "2":
                cmd += ["--two-stems", "vocals"]
            cmd.append(str(in_path))

            proc = await asyncio.create_subprocess_exec(
                *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await proc.communicate()
            if proc.returncode != 0:
                raise Exception(f"Demucs error: {stderr.decode()[-400:]}")

            out_dir = tmp / "htdemucs" / in_path.stem
            if not out_dir.exists():
                # fallback: sometimes model name differs
                candidates = list(tmp.glob("*/" + in_path.stem))
                out_dir = candidates[0] if candidates else out_dir

            JOBS[job_id]["progress"] = 80
            urls = []
            stem_names = ["vocals", "accompaniment"] if stems == "2" else ["vocals", "drums", "bass", "other"]
            for stem_name in stem_names:
                for ext in ("mp3", "wav"):
                    src = out_dir / f"{stem_name}.{ext}"
                    if src.exists():
                        dest = AUDIO_DIR / f"demucs-{job_id}-{stem_name}.{ext}"
                        shutil.copy(str(src), str(dest))
                        urls.append({"stem": stem_name, "url": f"/api/audio/{dest.name}", "filename": dest.name})
                        break

            if not urls:
                raise Exception("Demucs nie wygenerował plików wyjściowych — sprawdź logi")

            JOBS[job_id] = {
                "status": "completed", "progress": 100,
                "url": urls[0]["url"],   # primary (vocals) for backward compat
                "stems": urls,
                "message": f"Rozdzielono na {len(urls)} ścieżek",
            }
        except Exception as e:
            JOBS[job_id] = {"status": "failed", "error": str(e), "url": None, "stems": []}
        finally:
            shutil.rmtree(str(tmp), ignore_errors=True)

    background_tasks.add_task(_run)
    return {"job_id": job_id}

# ── FAZA 8: Security + Cache ───────────────────────────────────────────────────

_AUTH_CONFIG_PATH = Path(__file__).parent / "auth_config.json"

class _PasswordReq(BaseModel):
    password: str

@app.post("/api/auth/set-password")
async def set_password(req: _PasswordReq):
    try:
        import bcrypt as _bcrypt
    except ImportError:
        try:
            from passlib.hash import bcrypt as _bcrypt_passlib
            hashed = _bcrypt_passlib.hash(req.password)
            _AUTH_CONFIG_PATH.write_text(json.dumps({"password_hash": hashed}))
            return {"status": "password_set"}
        except ImportError:
            raise HTTPException(503, "Zainstaluj: pip install bcrypt lub pip install passlib[bcrypt]")
    hashed = _bcrypt.hashpw(req.password.encode(), _bcrypt.gensalt()).decode()
    _AUTH_CONFIG_PATH.write_text(json.dumps({"password_hash": hashed}))
    return {"status": "password_set"}

@app.get("/api/auth/check")
async def check_auth():
    has = _AUTH_CONFIG_PATH.exists()
    return {"password_enabled": has, "has_password": has}

@app.post("/api/cache/clear")
async def clear_cache():
    import time as _time
    tmp_dir = Path("/tmp")
    prefixes = ("vs_", "ve_", "ab_", "pres_", "music_", "demucs_")
    now = _time.time()
    cleared = 0
    freed_mb = 0.0
    for item in tmp_dir.iterdir():
        if not any(item.name.startswith(p) for p in prefixes):
            continue
        try:
            age = now - item.stat().st_mtime
            if age < 3600:
                continue
            size = sum(f.stat().st_size for f in item.rglob("*") if f.is_file()) if item.is_dir() else item.stat().st_size
            freed_mb += size / (1024 * 1024)
            if item.is_dir():
                shutil.rmtree(str(item), ignore_errors=True)
            else:
                item.unlink(missing_ok=True)
            cleared += 1
        except Exception:
            pass
    return {"cleared": cleared, "freed_mb": round(freed_mb, 2)}

# ──────────────────────────────────────────────────────────────────────────────
# Model Registry API
try:
    from model_registry import (
        MODEL_CATALOG, get_model_status_all, get_active_models,
        set_active_model, auto_configure_for_hardware, detect_hardware_profile,
    )
    _MODEL_REGISTRY_AVAILABLE = True
except ImportError:
    _MODEL_REGISTRY_AVAILABLE = False

class _ModelActivateRequest(BaseModel):
    category: str
    model_id: str

def _get_vram_total_gb() -> float:
    """Zwraca całkowitą ilość VRAM w GB. 0 jeśli brak GPU lub torch niedostępny."""
    try:
        import torch as _t
        if _t.cuda.is_available():
            return _t.cuda.get_device_properties(0).total_memory / 1e9
    except Exception:
        pass
    return 0.0

@app.get("/api/models")
async def list_models():
    if not _MODEL_REGISTRY_AVAILABLE:
        raise HTTPException(503, "Model registry niedostępny")
    vram = _get_vram_total_gb()
    return {
        "models": get_model_status_all(vram),
        "active": get_active_models(),
        "hardware_profile": detect_hardware_profile(vram),
        "vram_total_gb": round(vram, 1),
    }

@app.post("/api/models/active")
async def activate_model(request: _ModelActivateRequest):
    if not _MODEL_REGISTRY_AVAILABLE:
        raise HTTPException(503, "Model registry niedostępny")
    try:
        set_active_model(request.category, request.model_id)
        return {"status": "ok", "active": get_active_models()}
    except ValueError as e:
        raise HTTPException(400, str(e))

@app.post("/api/models/auto-configure")
async def auto_configure_models():
    if not _MODEL_REGISTRY_AVAILABLE:
        raise HTTPException(503, "Model registry niedostępny")
    vram = _get_vram_total_gb()
    return auto_configure_for_hardware(vram)

@app.get("/api/models/{model_id}/install-info")
async def get_model_install_info(model_id: str):
    if not _MODEL_REGISTRY_AVAILABLE:
        raise HTTPException(503, "Model registry niedostępny")
    model = MODEL_CATALOG.get(model_id)
    if not model:
        raise HTTPException(404, f"Model '{model_id}' nie istnieje w katalogu")
    return {
        "model_id": model_id,
        "name": model.get("name"),
        "script": model.get("install_script"),
        "cmd": model.get("install_cmd"),
        "disk_gb": model.get("disk_gb"),
        "vram_gb": model.get("vram_gb"),
    }

# ──────────────────────────────────────────────────────────────────────────────
# Static files for production
if FRONTEND_DIST.exists():
    app.mount("/", StaticFiles(directory=str(FRONTEND_DIST), html=True), name="static")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=47821)
